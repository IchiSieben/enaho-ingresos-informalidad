# -*- coding: utf-8 -*-
# verificar_ppt.py — reabre la presentación y la audita
# ---------------------------------------------------------------------------
# Autor: Yoichi Palacios Tanaka (IchiSieben) — proyecto ENAHO 2025
# Licencia: Apache-2.0 (ver LICENSE)
# ---------------------------------------------------------------------------
"""
Audita una salida de la exposición después de generarla (por defecto la
ENTREGA; pásale otra ruta como primer argumento para auditar la EXPO:
`verificar_ppt.py docs/presentacion/ENAHO_exposicion_EXPO.pptx`). Comprueba:

  1. estructura     — 13 láminas, cada una con notas del orador y con al menos
                      un elemento visual dominante (imagen o tabla);
  2. tipografía     — ninguna caja de texto con cuerpo por debajo de 18 pt,
                      salvo las dos excepciones declaradas (etiqueta de tarjeta
                      y pie de fuente), que no llevan prosa;
  3. cifras         — toda cifra visible cruzada contra el artefacto que la
                      genera, y tabla lámina → cifra → fuente.

Un `run` sin tamaño explícito hereda el del layout y NO se puede dar por bueno:
se cuenta como fallo, no como aprobado. Las tablas no exponen `text_frame`, hay
que entrar por `shape.has_table` — si no, la auditoría no ve ni una celda y
declara cero incumplimientos sobre una lámina entera de texto pequeño.

Correr:  .venv/Scripts/python.exe docs/presentacion/verificar_ppt.py
"""
from pathlib import Path
import csv
import json
import re
import sys

from pptx import Presentation
from pptx.util import Pt

RAIZ = Path(__file__).resolve().parents[2]
AQUI = RAIZ / "docs" / "presentacion"
PPTX = (Path(sys.argv[1]) if len(sys.argv) > 1
        else AQUI / "ENAHO_exposicion.pptx")
LAMINAS_ESPERADAS = 13
MINIMO_PT = 18.0          # cuerpo en cajas de texto
MINIMO_TABLA_PT = 16.0    # celdas de tabla: 5 columnas de prosa a 18 no caben

UA = json.loads((RAIZ / "models" / "ui_artifacts.json").read_text(encoding="utf-8"))
FS = json.loads((RAIZ / "models" / "feature_schema.json").read_text(encoding="utf-8"))
INFORME = (RAIZ / "INFORME_AUDITORIA.md").read_text(encoding="utf-8")
CONSOLA = (AQUI / "salida_consola_verificacion.txt").read_text(encoding="utf-8")
REQUISITOS = (RAIZ / "requirements.txt").read_text(encoding="utf-8")
sys.path.insert(0, str(RAIZ / "app"))
import referencias as REFS  # noqa: E402

fallos: list[str] = []
avisos: list[str] = []


def falla(msg):
    fallos.append(msg)
    print("  FALLA  " + msg)


def ok(msg):
    print("  ok     " + msg)


# ---------------------------------------------------------------------------
# 1. El diccionario de cifras permitidas, construido desde los artefactos
# ---------------------------------------------------------------------------
def _n(x):
    return f"{x:,.0f}".replace(",", ".")


def _d(x, dec=2):
    return f"{x:.{dec}f}".replace(".", ",")


PERMITIDAS: dict[str, str] = {}


def permitir(valor, fuente):
    """
    Registra una cifra y su procedencia en las formas en que PUEDE aparecer
    escrita en una lámina — y solo en esas.

    Generar todos los redondeos (`_d(x, 0)` para cualquier x) hacía que un
    «2» de la diapositiva casara con el smearing 1,564 redondeado, y la tabla
    de procedencia se llenaba de emparejamientos falsos que no prueban nada.
    """
    formas = {str(valor)}
    if isinstance(valor, (int, float)):
        v = float(valor)
        if v.is_integer():
            formas |= {str(int(v)), _n(v), _d(v, 1)}
        else:
            formas |= {_d(v, k) for k in (1, 2, 3, 4)}
            if abs(v) >= 100:                       # soles y conteos grandes
                formas |= {_n(v), _n(round(v))}
            if 0 < abs(v) < 1:                      # proporciones → porcentaje
                formas |= {_d(v * 100, k) for k in (0, 1, 2)}
    for f in formas:
        f = str(f).strip()
        if f and f not in PERMITIDAS:
            PERMITIDAS[f] = fuente


# --- del INFORME (mismos patrones que el generador) ---
for clave, pat in [
    ("embudo.crudo", r"Módulo 05 — Empleo e ingresos \(crudo\)\s+([\d.]+)"),
    ("embudo.ocupados", r"\nOcupados\s+([\d.]+)"),
    ("embudo.modelado", r"Dataset de modelado\s+([\d.]+) filas"),
    ("embudo.torneo", r"Muestra del torneo E1–E9\s+([\d.]+)"),
    ("embudo.train", r"Train ([\d.]+)"),
    ("embudo.test", r"Train [\d.]+\s+·\s+Test ([\d.]+)"),
    ("rejilla.vieja", r"3 de 3 en el borde\*?\*? \| ([\d]+,[\d]+) \|"),
    ("rejilla.nueva", r"ninguno en el borde \| \*\*([\d]+,[\d]+)\*\* \|"),
    ("rejilla.mejora_pct", r"−3,59 \(−([\d]+,[\d]+) %\)"),
    ("externa.propia", r"Nacional ponderado \| ([\d]+,[\d]+) %"),
    ("externa.inei", r"Nacional ponderado \| [\d,]+ % \| ([\d]+,[\d]+) %"),
    ("externa.inei_1_10", r"\| \*\*(88,6) %\*\* \| INEI, tramo \*\*1-10\*\*"),
]:
    m = re.search(pat, INFORME)
    if not m:
        falla(f"el patrón de {clave} ya no casa en INFORME_AUDITORIA.md")
    else:
        permitir(m.group(1), f"INFORME_AUDITORIA.md · {clave}")
        permitir(m.group(1).replace(".", ""), f"INFORME_AUDITORIA.md · {clave}")

for clave, pat, quien in [
    ("consola.ing_tipico", r"ingreso típico\s+\(mediana\): S/ ([\d.]+)", CONSOLA),
    ("consola.ing_esperado", r"ingreso esperado.*: +S/ ([\d.]+)", CONSOLA),
    ("consola.proba", r"probabilidad de informalidad: ([\d]+,[\d]+)%", CONSOLA),
    ("requirements.sklearn", r"scikit-learn==([\d.]+)", REQUISITOS),
    ("requirements.streamlit", r"streamlit==([\d.]+)", REQUISITOS),
    ("requirements.pandas", r"pandas==([\d.]+)", REQUISITOS),
]:
    m = re.search(pat, quien)
    if m:
        permitir(m.group(1), clave)
        permitir(m.group(1).replace(".", ""), clave)

# --- del torneo, el clasificador y el schema ---
for f in UA["torneo"]["tabla"]:
    for k in ("MAE_cv", "MAE_test", "MAE_test_media_smear", "RMSE_test",
              "R2_test_soles", "R2_escala_propia", "smearing_Duan", "n_vars"):
        if f.get(k) is not None:
            permitir(f[k], f"ui_artifacts · torneo.tabla[{f['ID']}].{k}")
for f in UA["clasificador"]["comparacion"]:
    for k in ("PRAUC_cv", "ROCAUC_cv", "PRAUC_test", "ROCAUC_test", "Brier_test"):
        permitir(f[k], f"ui_artifacts · clasificador.comparacion[{f['algoritmo']}].{k}")
for k, v in FS["clasificador"]["punto_operativo"].items():
    if isinstance(v, (int, float)):
        permitir(v, f"feature_schema · clasificador.punto_operativo.{k}")
for k in ("prevalencia_train", "prevalencia_ponderada", "n_entrenamiento", "n_test"):
    permitir(FS["clasificador"][k], f"feature_schema · clasificador.{k}")
for k in ("n_entrenamiento", "n_test", "smearing_duan", "ingreso_mediano_train"):
    permitir(FS["regresor"][k], f"feature_schema · regresor.{k}")
for k, v in FS["regresor"]["metricas_test"].items():
    permitir(v, f"feature_schema · regresor.metricas_test.{k}")
for k, v in FS["clasificador"]["metricas_test"].items():
    permitir(v, f"feature_schema · clasificador.metricas_test.{k}")
# Tasas observadas por grupo: el arreglo del hallazgo del 88,6 % fue
# justamente dejar de escribirlas a mano y calcularlas en el precómputo, así
# que la lámina de la auditoría las cita desde aquí.
for _var, _bloque in UA["clasificador"].get("tasas_observadas", {}).items():
    for _grupo, _datos in _bloque.get("grupos", {}).items():
        for _k in ("pct_ponderado", "pct_crudo"):
            if _k in _datos:
                permitir(_datos[_k],
                         f"ui_artifacts · tasas_observadas.{_var}[{_grupo}].{_k}")
permitir(UA["clasificador"]["pr"]["baseline"], "ui_artifacts · clasificador.pr.baseline")
permitir(UA["clasificador"]["pr"]["auc"], "ui_artifacts · clasificador.pr.auc")
permitir(UA["clasificador"]["roc"]["auc"], "ui_artifacts · clasificador.roc.auc")
permitir(UA["regresor"]["ingreso"]["mediana_ponderada"],
         "ui_artifacts · regresor.ingreso.mediana_ponderada")
for f in csv.DictReader(open(RAIZ / "reports" / "ablacion_clasificador.csv",
                             encoding="utf-8")):
    for k, v in f.items():
        if k != "variante":
            permitir(v, f"ablacion_clasificador.csv · {f['variante']}.{k}")
            try:
                permitir(float(v), f"ablacion_clasificador.csv · {f['variante']}.{k}")
            except ValueError:
                pass
imp = UA["regresor"]["importancia_permutacion"]
for nom, med in zip(imp["variables"], imp["media"]):
    permitir(round(med), f"ui_artifacts · regresor.importancia[{nom}]")
permitir(imp["n_filas"], "ui_artifacts · regresor.importancia.n_filas")
permitir(imp["n_repeticiones"], "ui_artifacts · regresor.importancia.n_repeticiones")
# La del clasificador se mide en puntos de PR-AUC (fracciones pequeñas):
# se permite el valor crudo, cuyas formas redondeadas (0,046…) son las que
# la lámina de variables escribe.
imp_c = UA["clasificador"]["importancia_permutacion"]
for nom, med in zip(imp_c["variables"], imp_c["media"]):
    permitir(med, f"ui_artifacts · clasificador.importancia[{nom}]")
permitir(imp_c["n_filas"], "ui_artifacts · clasificador.importancia.n_filas")
permitir(imp_c["n_repeticiones"],
         "ui_artifacts · clasificador.importancia.n_repeticiones")
permitir(len(REFS.REFERENCIAS), "app/referencias.py · nº de referencias")
permitir(len([f for f in FS["regresor"]["features"] if f["tipo"] == "numerico"]),
         "feature_schema · nº de numéricas")
permitir(len([f for f in FS["regresor"]["features"] if f["tipo"] == "categorico"]),
         "feature_schema · nº de categóricas")
permitir(len(FS["regresor"]["features"]), "feature_schema · nº de features")
for f in FS["regresor"]["features"]:
    for k in ("min", "max", "default"):
        if k in f:
            permitir(f[k], f"feature_schema · {f['nombre']}.{k}")
    if "opciones" in f:
        permitir(len(f["opciones"]), f"feature_schema · {f['nombre']}: nº de niveles")
permitir(sum(len(f["opciones"]) for f in FS["regresor"]["features"]
             if f["tipo"] == "categorico"), "feature_schema · niveles one-hot")
permitir(len(UA["clasificador"]["curva_umbral"]["umbral"]),
         "ui_artifacts · nº de puntos de la curva de umbral")
permitir(len(UA["clasificador"]["histograma_oof"]["clase_1"]),
         "ui_artifacts · nº de bins del histograma OOF")
permitir(len(UA["clasificador"]["dependencia_parcial"]),
         "ui_artifacts · nº de dependencias parciales")
permitir(UA["torneo"]["autopsia"]["n_centinelas"], "ui_artifacts · autopsia.n_centinelas")
permitir(UA["torneo"]["autopsia"]["pct_centinelas"], "ui_artifacts · autopsia.pct_centinelas")
for k, v in UA["torneo"]["autopsia"].get("corrida_sucia", {}).items():
    permitir(v, f"ui_artifacts · autopsia.corrida_sucia.{k}")
for k, v in UA["torneo"]["autopsia"].get("corrida_limpia", {}).items():
    permitir(v, f"ui_artifacts · autopsia.corrida_limpia.{k}")
permitir(UA["torneo"]["explicativo_e6_ponderado"]["r2"], "ui_artifacts · E6.r2")
for k, v in UA["torneo"]["explicativo_e6_ponderado"]["efectos_pct"].items():
    permitir(v, f"ui_artifacts · E6.efectos_pct.{k}")
for f in UA["torneo"]["sensibilidad_especie"]:
    for k, v in f.items():
        if isinstance(v, (int, float)):
            permitir(v, f"ui_artifacts · sensibilidad_especie.{k}")

# consecuencias del umbral, en las formas en que la app y la PPT las muestran
_c = UA["clasificador"]["curva_umbral"]
_N = _c["n"]
permitir(_N, "ui_artifacts · curva_umbral.n")
for k in range(len(_c["umbral"])):
    permitir(round(_c["tp"][k] * 1000 / _N), "ui_artifacts · curva_umbral por mil")
    permitir(round(_c["fp"][k] * 1000 / _N), "ui_artifacts · curva_umbral por mil")
    permitir(round(_c["tn"][k] * 1000 / _N), "ui_artifacts · curva_umbral por mil")
    permitir(round(_c["fn"][k] * 1000 / _N), "ui_artifacts · curva_umbral por mil")
    permitir(round(_c["precision_1"][k] * 100), "ui_artifacts · curva_umbral precisión")
    permitir(round(_c["recall_1"][k] * 100), "ui_artifacts · curva_umbral recall")
    permitir(round((_c["tp"][k] + _c["fp"][k]) / _N * 100),
             "ui_artifacts · curva_umbral cobertura")

# medidos del disco y del repositorio
import subprocess  # noqa: E402

_git = lambda *a: subprocess.run(["git", *a], cwd=RAIZ, capture_output=True,
                                 text=True, encoding="utf-8").stdout.strip()
_vers = [f for f in _git("ls-files").split("\n") if f]
permitir(len(_vers), "git ls-files · nº de archivos versionados")
permitir(int(_git("rev-list", "--count", "HEAD")), "git · nº de commits")
_SALIDAS = {"docs/presentacion/ENAHO_exposicion.pptx",
            "docs/presentacion/ENAHO_exposicion_EXPO.pptx"}
permitir(round(sum((RAIZ / f).stat().st_size for f in _vers
                   if (RAIZ / f).exists() and f not in _SALIDAS) / 1e6, 2),
         "disco · MB versionados (sin las propias presentaciones)")
_data = RAIZ / "data"
if _data.exists():
    permitir(round(sum(x.stat().st_size for x in _data.rglob("*") if x.is_file())
                   / 1e6, 1), "disco · MB de data/ (fuera del repo)")
permitir(round((RAIZ / "models" / "ui_artifacts.json").stat().st_size / 1024, 1),
         "disco · KB de ui_artifacts.json")
permitir(round(sum((RAIZ / "models" / x).stat().st_size for x in
                   ("regresor_e9.joblib", "clasificador_gb.joblib")) / 1024),
         "disco · KB de los dos .joblib")
for p in sorted((RAIZ / "src").glob("*.py")):
    m = re.match(r"^(\d+)_", p.name)
    if m:
        permitir(int(m.group(1)), f"src/ · numeración de scripts ({p.name})")
        permitir(m.group(1), f"src/ · numeración de scripts ({p.name})")

# Los centinelas del INEI son códigos, no mediciones, pero se escriben en la
# lámina de auditoría. Salen de la constante del pipeline, no de la memoria.
_cent = re.search(r"CENTINELAS_MONETARIOS = \[([\d.,\s]+)\]",
                  (RAIZ / "src" / "comun.py").read_text(encoding="utf-8"))
if _cent:
    for _v in _cent.group(1).split(","):
        permitir(_v.strip(), "src/comun.py · CENTINELAS_MONETARIOS")
        # También en su forma legible («999.999»): la lámina de la auditoría
        # lo escribe con separador de millar, porque ahí el centinela se
        # cuenta como el sueldo imposible que el modelo llegó a creerse.
        try:
            permitir(float(_v), "src/comun.py · CENTINELAS_MONETARIOS")
        except ValueError:
            pass
else:
    falla("no se encontró CENTINELAS_MONETARIOS en src/comun.py")

# Estructurales: no son cifras de resultado y no requieren artefacto.
# El 25 es el día de la exposición (carátula: 25 de agosto de 2026).
ESTRUCTURALES = {
    "2025", "2026", "2024", "2022", "16", "9", "1", "2", "3", "4", "5", "6", "7",
    "8", "10", "11", "12", "13", "14", "15", "17", "18", "0", "20", "25", "50",
    "100", "500", "1.000", "1000", "1,33", "0,90", "0,5", "0,50", "0,500",
    "2.0", "4.0",          # nombres de licencia: Apache-2.0, CC BY-NC 4.0
}
for e in ESTRUCTURALES:
    PERMITIDAS.setdefault(e, "estructural (numeración, formato, fecha)")

# ---------------------------------------------------------------------------
# 2. Recorrido de la presentación
# ---------------------------------------------------------------------------
if not PPTX.exists():
    raise SystemExit(f"No existe {PPTX}: corre antes generar_ppt.py")
prs = Presentation(PPTX)

print(f"\n=== 1. Estructura ({PPTX.name}) ===")
if len(prs.slides) == LAMINAS_ESPERADAS:
    ok(f"{len(prs.slides)} láminas")
else:
    falla(f"{len(prs.slides)} láminas, se esperaban {LAMINAS_ESPERADAS}")


def corridas(shape):
    """Todo run de texto del shape, incluidas las celdas de tabla.

    Las tablas NO exponen `text_frame`: sin esta rama la auditoría no ve ni
    una celda y da por buena una lámina entera de texto pequeño.
    """
    if shape.has_table:
        for fila in shape.table.rows:
            for celda in fila.cells:
                for p in celda.text_frame.paragraphs:
                    for r in p.runs:
                        yield r, "tabla"
        return
    if shape.shape_type == 6:                      # grupo
        for sub in shape.shapes:
            yield from corridas(sub)
        return
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            yield r, "texto"


resumen = []
for i, s in enumerate(prs.slides, 1):
    imgs = sum(1 for sh in s.shapes if sh.shape_type == 13)
    tabs = sum(1 for sh in s.shapes if sh.has_table)
    nota = (s.notes_slide.notes_text_frame.text or "").strip()
    frases = len([x for x in re.split(r"[.!?]\s", nota) if x.strip()])
    resumen.append((i, imgs, tabs, len(nota), frases))
    if not nota:
        falla(f"lámina {i}: sin notas del orador")
    elif frases < 3:
        avisos.append(f"lámina {i}: notas de solo {frases} frase(s)")
    if imgs == 0 and tabs == 0:
        falla(f"lámina {i}: no tiene elemento visual (ni imagen ni tabla)")

print(f"{'lám':>4} {'imgs':>5} {'tablas':>7} {'notas(car)':>11} {'frases':>7}")
for r in resumen:
    print(f"{r[0]:>4} {r[1]:>5} {r[2]:>7} {r[3]:>11} {r[4]:>7}")

# --- tipografía ---
print(f"\n=== 2. Tipografía (mínimo {MINIMO_PT:.0f} pt para cuerpo) ===")
EXCEPCIONES = []          # (lámina, texto, pt) declaradas y toleradas
INFRACCIONES = []
SIN_TAMANO = []
for i, s in enumerate(prs.slides, 1):
    for sh in s.shapes:
        for r, origen in corridas(sh):
            txt = (r.text or "").strip()
            if not txt:
                continue
            if r.font.size is None:
                SIN_TAMANO.append((i, txt[:48]))
                continue
            pt = r.font.size.pt
            if pt >= MINIMO_PT:
                continue
            es_pie = txt.startswith("Fuente:")
            es_etiqueta = (r.font.name or "") == "Consolas" and txt == txt.upper()
            es_tabla_ok = origen == "tabla" and pt >= MINIMO_TABLA_PT
            if es_pie or es_etiqueta or es_tabla_ok:
                EXCEPCIONES.append((i, txt[:44], pt, origen))
            else:
                INFRACCIONES.append((i, txt[:60], pt, origen))

if SIN_TAMANO:
    for i, t in SIN_TAMANO:
        falla(f"lámina {i}: run sin tamaño explícito (hereda del layout): {t!r}")
else:
    ok("todos los runs llevan tamaño explícito")
if INFRACCIONES:
    for i, t, pt, o in INFRACCIONES:
        falla(f"lámina {i}: cuerpo a {pt:g} pt ({o}) — {t!r}")
else:
    ok(f"ningún cuerpo por debajo de {MINIMO_PT:.0f} pt")
_tab = [e for e in EXCEPCIONES if e[3] == "tabla"]
_otras = [e for e in EXCEPCIONES if e[3] != "tabla"]
print(f"  excepciones declaradas: {len(_tab)} celdas de tabla a "
      f"{MINIMO_TABLA_PT:.0f} pt o más (5 columnas de prosa a 18 pt no caben) "
      f"y {len(_otras)} etiquetas de tarjeta y pies de fuente, que no llevan "
      f"prosa.")
_bajo = [e for e in _tab if e[2] < MINIMO_TABLA_PT]
if _bajo:
    for i, t, pt, _ in _bajo:
        falla(f"lámina {i}: celda de tabla a {pt:g} pt — {t!r}")

# --- cifras ---
print("\n=== 3. Cifras: lámina → cifra → fuente ===")
# Los dígitos pegados a letras no son cifras: son hashes de commit (`0c7a001`),
# nombres de script (`04_torneo`) o identificadores de especificación (`E9`).
# Sin los límites, el «001» de un hash entraba como cifra sin fuente.
NUMERO = re.compile(r"(?<![\w/])\d[\d.,]*(?![\w])")
filas, sin_fuente = [], []
for i, s in enumerate(prs.slides, 1):
    vistos = set()
    # Las notas del orador son lo que se DICE en voz alta: una cifra de ahí sin
    # artefacto detrás se afirma de memoria delante del público. Se auditan
    # igual que las cajas de texto — antes quedaban fuera del recorrido.
    trozos = [(r.text or "", "lámina") for sh in s.shapes
              for r, _ in corridas(sh)]
    trozos.append((s.notes_slide.notes_text_frame.text or "", "notas"))
    for texto, donde in trozos:
        for m in NUMERO.finditer(texto):
            tok = m.group(0).rstrip(".,")
            if not tok or (tok, donde) in vistos:
                continue
            vistos.add((tok, donde))
            fuente = PERMITIDAS.get(tok)
            if fuente:
                filas.append((i, tok if donde == "lámina" else tok + " ·notas",
                              fuente))
            else:
                ctx = texto[max(0, m.start() - 32):m.end() + 32].strip()
                sin_fuente.append((i, f"{tok} [{donde}]", ctx[:84]))

print(f"{'lám':>4}  {'cifra':<12} fuente")
for i, tok, f in filas:
    print(f"{i:>4}  {tok:<12} {f}")
print(f"\n  {len(filas)} cifras trazadas a su artefacto")
if sin_fuente:
    print("\n  SIN FUENTE (revisar una a una):")
    for i, tok, ctx in sin_fuente:
        falla(f"lámina {i}: «{tok}» no sale de ningún artefacto — en {ctx!r}")
else:
    ok("ninguna cifra sin artefacto que la respalde")

print("\n" + "=" * 72)
for a in avisos:
    print("  aviso: " + a)
if fallos:
    print(f"{len(fallos)} FALLOS")
    sys.exit(1)
print("La presentación pasa la auditoría.")
