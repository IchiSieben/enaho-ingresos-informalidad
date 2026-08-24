# -*- coding: utf-8 -*-
# generar_ppt.py — la exposición, generada desde los artefactos
# ---------------------------------------------------------------------------
# Autor: Yoichi Palacios Tanaka (IchiSieben) — proyecto ENAHO 2025
# Licencia: Apache-2.0 (ver LICENSE)
# ---------------------------------------------------------------------------
"""
Genera DOS salidas de la misma exposición (16:9, 14 láminas en el orden de
la pauta del docente: carátula · datos · variables · torneo · importancia ·
despliegue · verificación · auditoría · prácticas · cierre):

  - ENAHO_exposicion.pptx       (ENTREGA) — notas mínimas: solo la línea de
                                 fuente por lámina. Es la que se envía.
  - ENAHO_exposicion_EXPO.pptx  (EXPO)    — láminas idénticas byte a byte,
                                 con las notas completas del orador (QUÉ
                                 DIGO · TÉRMINOS · SI PREGUNTAN). Es la
                                 chuleta del equipo en la vista del
                                 presentador.

Flags:  --solo-entrega  |  --solo-expo   (sin flag se generan las dos).

Las láminas se CONSTRUYEN en un orden y se ENTREGAN en otro: el reorden se
hace sobre la lista de diapositivas al guardar (ORDEN_FINAL), no moviendo
código — así el generador conserva su estructura y el mazo sigue la pauta.

La mesa evalúa un curso de DESPLIEGUE de machine learning: el mazo defiende
la arquitectura y sus garantías; el detalle estadístico —torneo de nueve,
embudo, umbral, prevalencia— vive en las notas EXPO, listo para preguntas,
no en pantalla.

Reglas de composición (verificar_ppt.py comprueba las medibles):

  - cada título es una AFIRMACIÓN completa, nunca un tema;
  - cero rótulos de estructura: ni «Bloque», ni «lámina n de m»;
  - un elemento visual dominante por lámina;
  - franja de acento constante en la cabecera de todas las láminas;
  - capturas tratadas como capturas de producto: recorte útil, borde fino,
    esquinas suaves;
  - cuerpo a 18 pt mínimo en sans (Calibri); la monoespaciada se reserva para
    cifras, nombres de archivo y código. Excepciones declaradas: etiquetas de
    tarjeta (14 pt, versalita mono) y pies de fuente (12 pt), sin prosa.

Regla de oro, la misma de toda la auditoría: NINGÚN número escrito de memoria.
Todo sale de:

  - models/ui_artifacts.json          (torneo, autopsia, clasificador, meta)
  - models/feature_schema.json        (variables, targets, punto operativo)
  - reports/ablacion_clasificador.csv (ablación estructural)
  - INFORME_AUDITORIA.md              (embudo N, rejilla ampliada, INEI)
  - app/referencias.py                (n.º de referencias verificadas)
  - docs/presentacion/salida_consola_verificacion.txt (VS Code = Streamlit)
  - el propio repositorio y el disco    (peso versionado, tamaño de data/)

`buscar()` aborta la generación si un patrón no aparece: preferimos no tener
presentación a tener una con una cifra inventada.

Correr:  .venv/Scripts/python.exe docs/presentacion/generar_ppt.py
"""
from pathlib import Path
import csv
import json
import re
import sys

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

RAIZ = Path(__file__).resolve().parents[2]
AQUI = RAIZ / "docs" / "presentacion"
FIGS = AQUI / "figuras"
LOGOS = AQUI / "logos"
FIGS.mkdir(exist_ok=True)

# ---------------------------------------------------------------------------
# 1. Artefactos — la única fuente de cifras
# ---------------------------------------------------------------------------
UA = json.loads((RAIZ / "models" / "ui_artifacts.json").read_text(encoding="utf-8"))
FS = json.loads((RAIZ / "models" / "feature_schema.json").read_text(encoding="utf-8"))
INFORME = (RAIZ / "INFORME_AUDITORIA.md").read_text(encoding="utf-8")
CONSOLA = (AQUI / "salida_consola_verificacion.txt").read_text(encoding="utf-8")
REQUISITOS = (RAIZ / "requirements.txt").read_text(encoding="utf-8")

sys.path.insert(0, str(RAIZ / "app"))
import referencias as REFS  # noqa: E402

with open(RAIZ / "reports" / "ablacion_clasificador.csv", encoding="utf-8") as f:
    ABLACION = list(csv.DictReader(f))


def buscar(patron: str, texto: str = INFORME,
           quien: str = "INFORME_AUDITORIA.md") -> str:
    """Un número del informe se extrae, no se copia. Si no está, se aborta."""
    m = re.search(patron, texto)
    if not m:
        raise SystemExit(f"No se encontró {patron!r} en {quien}: no se inventa.")
    return m.group(1)


# Embudo de datos (INFORME §4, reconstruido reejecutando la Fase 1). En este
# mazo el embudo se CUENTA, no se proyecta: vive en las notas del orador.
EMBUDO = {
    "crudo": buscar(r"Módulo 05 — Empleo e ingresos \(crudo\)\s+([\d.]+)"),
    "ocupados": buscar(r"\nOcupados\s+([\d.]+)"),
    "modelado": buscar(r"Dataset de modelado\s+([\d.]+) filas"),
    "torneo": buscar(r"Muestra del torneo E1–E9\s+([\d.]+)"),
    "train": buscar(r"Train ([\d.]+)"),
    "test": buscar(r"Train [\d.]+\s+·\s+Test ([\d.]+)"),
}
# Rejilla ampliada del regresor (INFORME §6)
REJILLA = {
    "vieja": buscar(r"3 de 3 en el borde\*?\*? \| ([\d]+,[\d]+) \|"),
    "nueva": buscar(r"ninguno en el borde \| \*\*([\d]+,[\d]+)\*\* \|"),
    "mejora_pct": buscar(r"−3,59 \(−([\d]+,[\d]+) %\)"),
}
# Validación externa del target (INFORME §5) y AC-5 (el 88,6 %)
EXTERNA = {
    "propia": buscar(r"Nacional ponderado \| ([\d]+,[\d]+) %"),
    "inei": buscar(r"Nacional ponderado \| [\d,]+ % \| ([\d]+,[\d]+) %"),
    "inei_1_10": buscar(r"\| \*\*(88,6) %\*\* \| INEI, tramo \*\*1-10\*\*"),
}
N_REFERENCIAS = len(REFS.REFERENCIAS)
# La tasa propia de la categoría «Hasta 20», que es lo que el 88,6 % del INEI
# NO era. Sale del precómputo, que es exactamente el arreglo que se hizo:
# la cifra se calcula, no se escribe a mano (INFORME §2.2.1).
TASA_HASTA20 = (UA["clasificador"]["tasas_observadas"]["tamano_empresa"]
                ["grupos"]["Hasta 20"]["pct_ponderado"])

TAB_TORNEO = {f["ID"]: f for f in UA["torneo"]["tabla"]}
CLF_COMP = UA["clasificador"]["comparacion"]      # GB / RF / logística
REG, CLF = FS["regresor"], FS["clasificador"]
PUNTO = CLF["punto_operativo"]
META = UA["meta"]
CURVA = UA["clasificador"]["curva_umbral"]
BASE_PR = UA["clasificador"]["pr"]["baseline"]
IMP_REG = UA["regresor"]["importancia_permutacion"]
IMP_CLF = UA["clasificador"]["importancia_permutacion"]
E6 = UA["torneo"]["explicativo_e6_ponderado"]
ABL_V2 = next(f for f in ABLACION if "categoria" in f["variante"])

# La consola local y la app en la nube deben decir lo mismo — se verifica AQUÍ,
# al generar: si algún día divergen, la PPT no se genera.
ING_TIPICO = buscar(r"ingreso típico\s+\(mediana\): S/ ([\d.]+)", CONSOLA, "consola")
ING_ESPERADO = buscar(r"ingreso esperado.*: +S/ ([\d.]+)", CONSOLA, "consola")
PROBA = buscar(r"probabilidad de informalidad: ([\d]+,[\d]+%)", CONSOLA, "consola")
# La consola imprime «97,5%» pegado; en la lámina va con espacio, como en
# la app y como manda la ortografía.
PROBA_ES = PROBA.replace("%", " %")
SKLEARN_REQ = buscar(r"scikit-learn==([\d.]+)", REQUISITOS, "requirements.txt")
if SKLEARN_REQ != META["version_scikit_learn"]:
    raise SystemExit(
        f"requirements.txt fija scikit-learn {SKLEARN_REQ} pero los artefactos "
        f"se generaron con {META['version_scikit_learn']}: los pickles no "
        f"cargarían en la nube. No se genera la presentación.")

# Variables: las mismas 11 en los dos modelos. Se comprueba, no se supone.
NUM = [f for f in REG["features"] if f["tipo"] == "numerico"]
CAT = [f for f in REG["features"] if f["tipo"] == "categorico"]
if [f["nombre"] for f in REG["features"]] != [f["nombre"] for f in CLF["features"]]:
    raise SystemExit("Regresor y clasificador ya no comparten features: "
                     "la tabla de variables daría por buena una tabla que no lo es.")
N_NIVELES = sum(len(f["opciones"]) for f in CAT)

# Peso de lo versionado y de lo que se queda fuera — medido, no recordado.
# El número de commits NO se imprime en ninguna lámina: el commit que publica
# la presentación lo cambiaría y la cifra nacería desactualizada.
import subprocess  # noqa: E402


def _git(*a) -> str:
    return subprocess.run(["git", *a], cwd=RAIZ, capture_output=True, text=True,
                          encoding="utf-8").stdout.strip()


SALIDAS_REL = {"docs/presentacion/ENAHO_exposicion.pptx",
               "docs/presentacion/ENAHO_exposicion_EXPO.pptx"}
VERSIONADOS = [f for f in _git("ls-files").split("\n") if f]
REPO = {
    "n_archivos": len(VERSIONADOS),
    # Sin contar las propias presentaciones (ENTREGA y EXPO): si entraran,
    # la cifra que la lámina de arquitectura imprime cambiaría al guardarlas
    # y no habría forma de verificarla (los archivos se miden antes de
    # existir en su tamaño final).
    "mb": round(sum((RAIZ / f).stat().st_size for f in VERSIONADOS
                    if (RAIZ / f).exists() and f not in SALIDAS_REL) / 1e6, 2),
}
_DATA = RAIZ / "data"
DATOS_MB = (round(sum(x.stat().st_size for x in _DATA.rglob("*") if x.is_file())
                  / 1e6, 1) if _DATA.exists() else None)
UI_KB = round((RAIZ / "models" / "ui_artifacts.json").stat().st_size / 1024, 1)
JOBLIB_KB = round(sum((RAIZ / "models" / n).stat().st_size for n in
                      ("regresor_e9.joblib", "clasificador_gb.joblib")) / 1024)

# ---------------------------------------------------------------------------
# 2. Paleta y tipografía — las del tema claro de la app (app/estilos.py)
# ---------------------------------------------------------------------------
C = {
    "fondo": "F7F5F0", "superficie": "FCFBF8", "superficie_alta": "F1EFE8",
    "borde": "D5D1C6", "texto": "1E232B", "texto_medio": "49525F",
    "texto_tenue": "5F6875", "acento": "4353CC", "acento_alto": "3542B8",
    "acento_fondo": "E6E8FA", "buena": "177A4C", "buena_fondo": "E2F2E9",
    "media": "8A5D0B", "media_fondo": "F6ECD6", "mala": "B92F33",
    "mala_fondo": "F9E3E4", "consola_fondo": "11151C", "consola_texto": "D6E2D0",
}
FUENTE = "Calibri"
MONO = "Consolas"

# Escala tipográfica. El mínimo de cuerpo es 18: proyectado, menos no se lee.
T_TITULO, T_ENTRADILLA, T_CUERPO = 36, 20, 18
T_CIFRA, T_CIFRA_XL = 34, 54
T_TABLA, T_TABLA_CAB = 18, 16
# Excepciones DECLARADAS al mínimo de 18 pt: ninguna lleva prosa.
T_ETIQUETA, T_PIE = 14, 12


def rgb(clave):
    return RGBColor.from_string(C[clave])


def mpl(clave):
    return "#" + C[clave]


def n(x) -> str:
    """1234.5 -> '1.234' (formato español, como la app)."""
    return f"{x:,.0f}".replace(",", ".")


def d(x, dec=2) -> str:
    return f"{x:.{dec}f}".replace(".", ",")


# ---------------------------------------------------------------------------
# 3. Figuras (matplotlib). Cada una es el elemento dominante de su lámina.
# ---------------------------------------------------------------------------
plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": [FUENTE, "Segoe UI", "DejaVu Sans"],
    "figure.facecolor": mpl("fondo"), "axes.facecolor": mpl("fondo"),
    "savefig.facecolor": mpl("fondo"), "text.color": mpl("texto"),
    "axes.labelcolor": mpl("texto_medio"), "xtick.color": mpl("texto_medio"),
    "ytick.color": mpl("texto_medio"), "axes.edgecolor": mpl("borde"),
})


def _guardar(fig, nombre):
    p = FIGS / f"{nombre}.png"
    fig.savefig(p, dpi=200, bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)
    return p


def _limpiar(ax, ejes=("top", "right")):
    for e in ejes:
        ax.spines[e].set_visible(False)


def fig_arquitectura():
    # Aspecto ancho a propósito: la lámina la encaja por altura y así el
    # diagrama ocupa casi todo el ancho útil en vez de quedar encogido.
    fig, ax = plt.subplots(figsize=(14.2, 5.1))
    ax.set_xlim(0, 100); ax.set_ylim(0, 50); ax.axis("off")

    def caja(x, y, w, h, titulo, lineas, fc, ec, tc=None):
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.7",
                                    fc=fc, ec=ec, lw=1.8, zorder=3))
        ax.text(x + w / 2, y + h - 2.6, titulo, ha="center", va="top",
                fontsize=16.5, fontweight="bold", color=tc or mpl("texto"),
                zorder=4)
        ax.text(x + w / 2, y + h - 7.6, "\n".join(lineas), ha="center", va="top",
                fontsize=13, color=mpl("texto_medio"), linespacing=1.6, zorder=4)

    def flecha(x1, y1, x2, y2, texto=None, dx=0, dy=2.0, ha="center"):
        ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>",
                                     mutation_scale=26, lw=2.6, color=mpl("acento"),
                                     zorder=5, shrinkA=0, shrinkB=0))
        if texto:
            ax.text((x1 + x2) / 2 + dx, (y1 + y2) / 2 + dy, texto, ha=ha,
                    va="center", fontsize=13.5, color=mpl("acento_alto"),
                    fontweight="bold", zorder=6,
                    bbox=dict(boxstyle="round,pad=0.22", fc=mpl("fondo"), ec="none"))

    caja(1, 28, 24, 19, "1 · Microdatos INEI",
         ["data/  ·  local", "ENAHO 2025", "módulos 02 · 03 · 05"],
         mpl("superficie_alta"), mpl("borde"))
    caja(31.5, 28, 24, 19, "2 · VS Code",
         ["src/00 → src/09", "torneo, entrenamiento", "y precómputo"],
         mpl("superficie"), mpl("borde"))
    caja(62, 28, 24, 19, "3 · Artefactos",
         ["2 modelos .joblib", "ui_artifacts.json", "feature_schema.json"],
         mpl("acento_fondo"), mpl("acento"))
    flecha(25.4, 37.5, 31.1, 37.5, "lee")
    flecha(55.9, 37.5, 61.6, 37.5, "escribe")
    flecha(74, 27.6, 74, 20.4, "git push", dx=5.6, dy=0, ha="left")
    # Fila inferior a 17,5 de alto, no 16: con el cuerpo a 13 pt la tercera
    # línea se salía por debajo del borde de la caja («.streamlit.app»
    # cortado). El alto se calcula para el texto, no al revés.
    caja(62, 3, 24, 17.5, "4 · GitHub",
         ["IchiSieben /", "enaho-ingresos-informalidad", "rama main"],
         mpl("superficie"), mpl("borde"))
    caja(31.5, 3, 24, 17.5, "5 · Streamlit Cloud",
         ["app/streamlit_app.py", "redeploy automático", "en cada push"],
         mpl("superficie"), mpl("borde"))
    caja(1, 3, 24, 17.5, "6 · Navegador",
         ["enaho-ingresos-informalidad", ".streamlit.app"],
         mpl("buena_fondo"), mpl("buena"), tc=mpl("buena"))
    flecha(61.6, 11.75, 55.9, 11.75, "webhook")
    flecha(31.1, 11.75, 25.4, 11.75, "HTTPS")
    ax.add_patch(FancyBboxPatch((29.6, 26.4), 58, 22.2, boxstyle="round,pad=0.5",
                                fc="none", ec=mpl("acento"), lw=1.4,
                                ls=(0, (6, 4)), zorder=2))
    ax.text(58.8, 49.6,
            f"carpeta versionada · {d(REPO['mb'], 2)} MB · "
            f"{REPO['n_archivos']} archivos",
            ha="center", va="bottom", fontsize=14.5, color=mpl("acento_alto"),
            fontweight="bold")
    # El mensaje de los microdatos NO va aquí dentro: incrustado en el hueco
    # entre las dos filas de cajas flotaba suelto y descentraba la figura al
    # recortarla `bbox_inches="tight"`. Vive como banda de ancho completo
    # bajo el diagrama, en la propia lámina.
    return _guardar(fig, "fig_arquitectura")


def fig_precomputo():
    # Aspecto 3,9:1: la lámina encaja la figura por altura (3,10") y con
    # este aspecto el resultado llena el ancho útil completo (12,09") en
    # vez de quedar encogido al centro con aire a los lados.
    fig, ax = plt.subplots(figsize=(15.2, 3.9))
    ax.set_xlim(0, 100); ax.set_ylim(0, 40); ax.axis("off")

    def bloque(x, w, titulo, subt, items, fc, ec, tc):
        """
        `items` son pares (qué es en llano, nombre técnico). El nombre
        técnico va debajo, gris y pequeño: quien lo conoce lo reconoce y
        quien no, ya entendió la línea de arriba.
        """
        ax.add_patch(FancyBboxPatch((x, 4), w, 32, boxstyle="round,pad=0.6",
                                    fc=fc, ec=ec, lw=1.6, zorder=3))
        ax.text(x + w / 2, 33.4, titulo, ha="center", va="top", fontsize=18.5,
                fontweight="bold", color=tc, zorder=4)
        ax.text(x + w / 2, 29.6, subt, ha="center", va="top", fontsize=13.5,
                color=mpl("texto_medio"), zorder=4)
        for k, (llano, tecnico) in enumerate(items):
            # Arranque 26,0 y paso 4,25: con el paso anterior (4,5) la
            # acotación del quinto ítem quedaba cortada por el borde
            # inferior de la caja (y=4) — el QA lo cazó en el render.
            cy = 26.0 - k * 4.25
            ax.text(x + 2.2, cy, "·  " + llano, ha="left", va="top",
                    fontsize=14.5, color=mpl("texto"), zorder=4)
            ax.text(x + 3.7, cy - 2.2, tecnico, ha="left", va="top",
                    fontsize=11.8, color=mpl("texto_tenue"), zorder=4)

    n_umbral = len(CURVA["umbral"])
    n_bins = len(UA["clasificador"]["histograma_oof"]["clase_1"])
    n_pdp = len(UA["clasificador"]["dependencia_parcial"])
    # Las cajas llegan casi al borde del lienzo (0,8..99,2): la lámina las
    # encaja por ancho y así el diagrama alcanza de verdad el ancho útil,
    # alineado con la banda del recorrido de abajo.
    bloque(0.8, 47.5, "UNA VEZ, en tu máquina", "src/09_precomputar_ui.py",
           [("Qué pasaría con cada corte del umbral",
             f"(curva de umbral · {n_umbral} puntos)"),
            ("Cómo se reparten las probabilidades del modelo",
             f"(histograma fuera de muestra · {n_bins} tramos)"),
            ("Cómo cambia la predicción con cada variable",
             f"(dependencia parcial · {n_pdp} variables)"),
            ("Con qué trabajadores parecidos se compara el perfil",
             "(cohortes y percentiles)"),
            ("Qué porcentaje es informal en cada grupo",
             "(tasas observadas)")],
           mpl("acento_fondo"), mpl("acento"), mpl("acento_alto"))
    bloque(51.7, 47.5, "EN CALIENTE, al pulsar el botón", "app/streamlit_app.py",
           [("Estimar el ingreso del perfil",
             "(una llamada a .predict() del regresor)"),
            ("Estimar su probabilidad de informalidad",
             "(una llamada a .predict_proba() del clasificador)"),
            ("Lo demás: leer el archivo guardado y dibujar",
             "(JSON → SVG)")],
           mpl("buena_fondo"), mpl("buena"), mpl("buena"))
    ax.annotate("", xy=(51.4, 20), xytext=(48.6, 20),
                arrowprops=dict(arrowstyle="-|>", mutation_scale=22, lw=2.2,
                                color=mpl("texto_tenue")))
    ax.text(24, 1.3, f"ui_artifacts.json · {d(UI_KB, 1)} KB", ha="center",
            va="center", fontsize=15.5, fontweight="bold",
            color=mpl("acento_alto"))
    ax.text(76, 1.3, f"2 modelos .joblib · {n(JOBLIB_KB)} KB", ha="center",
            va="center", fontsize=15.5, fontweight="bold", color=mpl("buena"))
    return _guardar(fig, "fig_precomputo")


def fig_consola(bloque: str, nombre: str):
    """
    La salida REAL de verificacion_local.py, pintada como una terminal.

    Solo la línea de cabecera y el bloque del modelo pedido: el perfil de
    entrada (los once valores por defecto) se cuenta en las notas, no se
    proyecta. Así las dos consolas caben en una única lámina de verificación.
    """
    lineas = [l.rstrip() for l in CONSOLA.strip().split("\n")]
    inicio = next(k for k, l in enumerate(lineas) if l.startswith(f"[{bloque}"))
    fin = next((k for k, l in enumerate(lineas[inicio + 1:], inicio + 1)
                if l.startswith("[")), len(lineas))
    lineas = [lineas[0]] + [l for l in lineas[inicio:fin] if l.strip()]
    # El ancho se calcula del texto, no se tantea: con `bbox_inches="tight"` un
    # renglón que se sale ensancha el lienzo y deja el panel oscuro corto, y la
    # línea más larga aparece pintada sobre el fondo claro.
    PT_CAR, MARGEN = 0.0965, 0.22          # pulgadas por carácter a 11,5 pt
    ancho = 2 * MARGEN + max(len(l) for l in lineas) * PT_CAR
    alto = 0.285 * len(lineas) + 0.5
    fig = plt.figure(figsize=(ancho, alto), facecolor=mpl("consola_fondo"))
    ax = fig.add_axes((0, 0, 1, 1))
    ax.axis("off")
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    for k, l in enumerate(lineas):
        resalta = ("S/ " in l or "%" in l or "veredicto" in l)
        ax.text(MARGEN / ancho, 1 - (k + 0.9) / (len(lineas) + 1.0), l,
                family="monospace", fontsize=11.5,
                color="#FFD43B" if resalta else mpl("consola_texto"),
                fontweight="bold" if resalta else "normal", va="center", zorder=2)
    p = FIGS / f"{nombre}.png"
    fig.savefig(p, dpi=200, facecolor=mpl("consola_fondo"))
    plt.close(fig)
    return p


def fig_requisitos():
    """requirements.txt frente a la versión que realmente generó los pickles."""
    lineas = [l for l in REQUISITOS.strip().splitlines()
              if l.strip() and not l.startswith("#")]
    fig = plt.figure(figsize=(11.4, 4.9), facecolor=mpl("fondo"))
    ax = fig.add_axes((0, 0, 1, 1)); ax.axis("off")
    ax.set_xlim(0, 100); ax.set_ylim(-7, 40)
    ax.add_patch(FancyBboxPatch((1, 2), 46, 36, boxstyle="round,pad=0.6",
                                fc=mpl("consola_fondo"), ec=mpl("consola_fondo")))
    ax.text(4, 35.6, "requirements.txt", family="monospace", fontsize=13,
            color="#FFD43B", fontweight="bold", va="top")
    for k, l in enumerate(lineas):
        clave = l.startswith("scikit-learn")
        ax.text(4, 31.4 - k * 2.62, l, family="monospace", fontsize=12.5,
                color="#FFD43B" if clave else mpl("consola_texto"),
                fontweight="bold" if clave else "normal", va="top")
    ax.add_patch(FancyBboxPatch((53, 2), 46, 36, boxstyle="round,pad=0.6",
                                fc=mpl("acento_fondo"), ec=mpl("acento"), lw=1.8))
    ax.text(76, 35.6, "models/ui_artifacts.json  ·  meta", ha="center", va="top",
            family="monospace", fontsize=13, color=mpl("acento_alto"),
            fontweight="bold")
    filas = [("version_scikit_learn", META["version_scikit_learn"]),
             ("version_pandas", META["version_pandas"]),
             ("fecha_generacion", META["fecha_generacion"]),
             ("commit", META["commit"][:7])]
    for k, (kk, vv) in enumerate(filas):
        ax.text(56, 29.4 - k * 4.4, kk, family="monospace", fontsize=12.5,
                color=mpl("texto_medio"), va="top")
        ax.text(96, 29.4 - k * 4.4, str(vv), family="monospace", fontsize=12.5,
                color=mpl("texto"), fontweight="bold", va="top", ha="right")
    ax.text(50, 20, "=", ha="center", va="center", fontsize=34,
            color=mpl("buena"), fontweight="bold")
    # El mensaje central de la lámina, no una letra chica: a tamaño de
    # subtítulo (la figura se coloca a ~0,73× en la lámina, así que estos
    # 27 pt quedan en ~20 pt proyectados), centrado bajo el «=».
    ax.text(50, -4.6, "si no coinciden, el modelo no carga en la nube",
            ha="center", va="center", fontsize=27, color=mpl("mala"),
            fontweight="bold")
    p = FIGS / "fig_requisitos.png"
    fig.savefig(p, dpi=200, facecolor=mpl("fondo"))
    plt.close(fig)
    return p


def fig_importancia():
    eti = {f["nombre"]: f["etiqueta"] for f in REG["features"]}
    # Las del schema son etiquetas de formulario: en un eje no caben enteras.
    eti.update({"exper2": "Experiencia² (cuadrado)",
                "horas_total": "Horas trabajadas por semana",
                "tamano_empresa": "Tamaño de la empresa",
                "anios_educ": "Años de educación aprobados"})
    it = sorted(zip(IMP_REG["variables"], IMP_REG["media"], IMP_REG["desviacion"]),
                key=lambda t: -t[1])
    fig, ax = plt.subplots(figsize=(10.8, 6.0))
    y = list(range(len(it)))[::-1]
    for k, (nom, med, des) in enumerate(it):
        ax.barh(y[k], med, height=0.6,
                color=mpl("acento") if k < 2 else mpl("texto_tenue"),
                alpha=1 if k < 2 else 0.55, xerr=des,
                error_kw=dict(ecolor=mpl("texto_medio"), lw=1.2, capsize=3),
                zorder=3)
        ax.text(med + des + 2.5, y[k], f"S/ {n(med)}", va="center", fontsize=14,
                fontweight="bold", color=mpl("texto"))
        ax.text(-2.5, y[k], eti.get(nom, nom), va="center", ha="right",
                fontsize=13,
                color=mpl("texto") if k < 2 else mpl("texto_medio"))
    ax.set_xlim(0, max(m for _, m, _ in it) * 1.22)
    ax.set_ylim(-0.7, len(it) - 0.3); ax.set_yticks([])
    ax.set_xlabel("Soles/mes de error añadido al barajar la variable al azar\n"
                  f"(permutación, {IMP_REG['n_repeticiones']} repeticiones sobre "
                  f"{n(IMP_REG['n_filas'])} filas de test)",
                  fontsize=12.5, linespacing=1.5)
    ax.tick_params(axis="x", labelsize=11.5)
    ax.grid(axis="x", color="#E4E1D8", lw=0.9, zorder=0)
    _limpiar(ax, ("top", "right", "left"))
    return _guardar(fig, "fig_importancia")


def fig_una_fuente():
    """
    El principio estrella en un vistazo: UNA fuente con tres lectores —
    y al lado, tachado, el antipatrón de las dos copias. Dos colores:
    acento para el patrón, rojo para lo que no se hace.
    """
    fig, ax = plt.subplots(figsize=(14.5, 1.75))
    ax.set_xlim(0, 100); ax.set_ylim(0, 12); ax.axis("off")

    def cajita(x, y0, w, h, texto, fc, ec, tc, fs=12.5, bold=False):
        ax.add_patch(FancyBboxPatch((x, y0), w, h, boxstyle="round,pad=0.35",
                                    fc=fc, ec=ec, lw=1.5, zorder=3))
        ax.text(x + w / 2, y0 + h / 2, texto, ha="center", va="center",
                fontsize=fs, fontweight="bold" if bold else "normal",
                color=tc, zorder=4)

    cajita(1, 3.4, 15, 5.2, "UNA FUENTE", mpl("acento_fondo"), mpl("acento"),
           mpl("acento_alto"), fs=14, bold=True)
    for _yc, _destino in ((10.0, "formulario"), (6.0, "este PPT"),
                          (2.0, "nube")):
        cajita(30, _yc - 1.7, 13.5, 3.4, _destino, mpl("superficie"),
               mpl("borde"), mpl("texto"))
        ax.add_patch(FancyArrowPatch((16.4, 6), (29.6, _yc),
                                     arrowstyle="-|>", mutation_scale=16,
                                     lw=1.8, color=mpl("acento"), zorder=2,
                                     shrinkA=0, shrinkB=0))
    ax.plot([58, 58], [0.8, 11.2], color=mpl("borde"), lw=1.2,
            ls=(0, (4, 3)), zorder=1)
    cajita(64, 4.2, 12.5, 3.6, "copia A", mpl("mala_fondo"), mpl("mala"),
           mpl("mala"))
    cajita(85, 4.2, 12.5, 3.6, "copia B", mpl("mala_fondo"), mpl("mala"),
           mpl("mala"))
    ax.text(80.8, 6.0, "≠", ha="center", va="center", fontsize=17,
            fontweight="bold", color=mpl("mala"), zorder=4)
    ax.text(80.8, 10.3, "?", ha="center", va="center", fontsize=15,
            fontweight="bold", color=mpl("mala"), zorder=4)
    # El tachado: este lado del diagrama es lo que NO se hace.
    ax.plot([62.5, 99], [0.6, 11.4], color=mpl("mala"), lw=2.4, zorder=5,
            alpha=0.85)
    return _guardar(fig, "fig_una_fuente")


F_REQ = fig_requisitos()
F_FUENTE = fig_una_fuente()
F_CONSOLA_REG = fig_consola("REGRESOR", "fig_consola_regresor")
F_CONSOLA_CLF = fig_consola("CLASIFICADOR", "fig_consola_clasificador")
F_ARQ = fig_arquitectura()
F_IMP = fig_importancia()
F_PRECOMP = fig_precomputo()


# ---------------------------------------------------------------------------
# 4. Utilidades python-pptx
# ---------------------------------------------------------------------------
ANCHO, ALTO = Inches(13.333), Inches(7.5)
MARGEN = 0.62                 # margen izquierdo/derecho común
UTIL = 13.333 - 2 * MARGEN    # ancho útil
prs = Presentation()
prs.slide_width, prs.slide_height = ANCHO, ALTO
BLANCO = prs.slide_layouts[6]

TITULOS: list[str] = []       # se imprimen al final: el mazo en una lista


def lamina(fondo="fondo"):
    s = prs.slides.add_slide(BLANCO)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(fondo)
    # La franja de acento constante: la misma identidad en todas las láminas.
    fr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ANCHO, Inches(0.07))
    fr.fill.solid()
    fr.fill.fore_color.rgb = rgb("acento")
    fr.line.fill.background()
    fr.shadow.inherit = False
    return s


def caja(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def parrafo_mixto(tf, partes, primero=False, alin=None, esp_despues=8,
                  esp_linea=None):
    """
    Un párrafo con varios tramos de formato distinto, en la MISMA línea base.

    `partes` son tuplas (texto, tamaño, negrita, color, fuente). Sirve para
    lo que un párrafo por tramo no puede dar: «App:» en sans junto a su URL
    en mono sin que la etiqueta caiga en su propio renglón.
    """
    p = (tf.paragraphs[0] if primero and not tf.paragraphs[0].runs
         else tf.add_paragraph())
    for texto, tam, negrita, color, fuente in partes:
        r = p.add_run()
        r.text = texto
        r.font.name = fuente
        r.font.size = Pt(tam)
        r.font.bold = negrita
        r.font.color.rgb = rgb(color)
    if alin:
        p.alignment = alin
    p.space_after = Pt(esp_despues)
    if esp_linea:
        p.line_spacing = esp_linea
    return p


def parrafo(tf, texto, tam=T_CUERPO, color="texto", negrita=False, fuente=FUENTE,
            primero=False, alin=None, esp_despues=8, esp_linea=None):
    p = (tf.paragraphs[0] if primero and not tf.paragraphs[0].runs
         else tf.add_paragraph())
    r = p.add_run()
    r.text = texto
    r.font.name = fuente
    r.font.size = Pt(tam)
    r.font.bold = negrita
    r.font.color.rgb = rgb(color)
    if alin:
        p.alignment = alin
    p.space_after = Pt(esp_despues)
    if esp_linea:
        p.line_spacing = esp_linea
    return p


# Rejilla fija: todas las láminas empiezan el cuerpo a la misma altura. Un
# título de una línea deja aire; uno de dos lo ocupa. Lo que NO puede pasar es
# que el contenido baile de lámina en lámina.
Y_TITULO, Y_ENTRADILLA, Y_CUERPO, Y_PIE = 0.42, 1.54, 2.26, 7.02
ALTO_CUERPO = Y_PIE - Y_CUERPO - 0.16


def titulo(s, texto, entradilla=""):
    """Título de lámina (1-2 líneas a 36 pt) más una línea llana debajo."""
    TITULOS.append(texto)
    tf = caja(s, MARGEN, Y_TITULO, UTIL, 1.06)
    parrafo(tf, texto, tam=T_TITULO, negrita=True, primero=True, esp_despues=0,
            esp_linea=0.92)
    if entradilla:
        tfe = caja(s, MARGEN, Y_ENTRADILLA, min(UTIL, 11.9), 0.66)
        parrafo(tfe, entradilla, tam=T_ENTRADILLA, color="texto_medio",
                primero=True, esp_despues=0, esp_linea=1.02)
    return Y_CUERPO


def panel(s, x, y, w, h, relleno="superficie", borde="borde"):
    fig = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    fig.adjustments[0] = 0.05
    fig.fill.solid()
    fig.fill.fore_color.rgb = rgb(relleno)
    fig.line.color.rgb = rgb(borde)
    fig.line.width = Pt(1.25)
    fig.shadow.inherit = False
    tf = fig.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.20)
    tf.margin_top = tf.margin_bottom = Inches(0.14)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return fig, tf


def tarjeta(s, x, y, w, h, etiqueta, cifra, nota, color_cifra="acento",
            relleno="superficie", tam_cifra=T_CIFRA):
    """
    Etiqueta + cifra + nota. Las tres son obligatorias: una tarjeta sin nota
    es una caja con aire, y de esas hemos aprendido a desconfiar.
    """
    _, tf = panel(s, x, y, w, h, relleno=relleno)
    parrafo(tf, etiqueta.upper(), tam=T_ETIQUETA, color="texto_tenue",
            fuente=MONO, primero=True, esp_despues=4)
    parrafo(tf, cifra, tam=tam_cifra, color=color_cifra, negrita=True,
            fuente=MONO, esp_despues=6, esp_linea=0.95)
    parrafo(tf, nota, tam=T_CUERPO, color="texto_medio", esp_despues=0,
            esp_linea=1.05)


def tabla(s, x, y, w, h, filas, anchos=None, tam=T_TABLA, tam_cab=T_TABLA_CAB,
          resaltar=None, resaltar_col=None, cols_mono=None):
    """
    `cols_mono` son las columnas de cifra: monoespaciada y alineadas a la
    derecha. Las demás son prosa, en sans y a la izquierda. `resaltar` marca
    una fila (el ganador de una comparación por filas); `resaltar_col`, una
    columna (el ganador cuando los competidores son las columnas).
    """
    cols_mono = (set(range(1, len(filas[0]))) if cols_mono is None
                 else set(cols_mono))
    nf, nc = len(filas), len(filas[0])
    gt = s.shapes.add_table(nf, nc, Inches(x), Inches(y), Inches(w),
                            Inches(h)).table
    if anchos:
        total = sum(anchos)
        for j, a in enumerate(anchos):
            gt.columns[j].width = Emu(int(Inches(w) * a / total))
    for i, fila in enumerate(filas):
        for j, valor in enumerate(fila):
            celda = gt.cell(i, j)
            celda.fill.solid()
            destacada = ((resaltar is not None and i == resaltar) or
                         (resaltar_col is not None and j == resaltar_col and i > 0))
            if i == 0:
                celda.fill.fore_color.rgb = rgb("superficie_alta")
            elif destacada:
                celda.fill.fore_color.rgb = rgb("acento_fondo")
            else:
                celda.fill.fore_color.rgb = rgb("superficie")
            celda.margin_left = celda.margin_right = Inches(0.10)
            celda.margin_top = celda.margin_bottom = Inches(0.05)
            celda.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = celda.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(valor)
            r.font.size = Pt(tam_cab if i == 0 else tam)
            r.font.name = MONO if (i > 0 and j in cols_mono) else FUENTE
            r.font.bold = (i == 0) or destacada
            r.font.color.rgb = rgb("texto_tenue" if i == 0 else "texto")
            if j in cols_mono and i > 0:
                p.alignment = PP_ALIGN.RIGHT
    return gt


def imagen(s, ruta, x, y, w=None, h=None):
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    return s.shapes.add_picture(str(ruta), Inches(x), Inches(y), **kw)


def recorte_lateral(nombre: str, destino: str, fraccion: float = 0.205):
    """
    Quita la barra lateral de una captura de la app y guarda el recorte.

    Se hace AQUÍ y no a mano: si el recorte viviera en un script suelto, el
    repositorio tendría una figura que nadie sabe reproducir. Lo que se
    versiona es la captura entera; el recorte se regenera en cada corrida.
    """
    from PIL import Image as _Img
    with _Img.open(FIGS / nombre) as im:
        an, al = im.size
        salida = FIGS / destino
        im.crop((int(an * fraccion), 0, an, al)).save(salida)
    return salida


def recorte_caja(nombre: str, destino: str, caja: tuple):
    """
    Recorta una zona (fracciones izq, arriba, der, abajo) de una captura.

    Para la lámina de verificación: de la página completa de la app solo se
    proyecta la zona donde vive el número que se compara — proyectar la página
    entera encogida es proyectar nada. Las fracciones están calibradas sobre
    las capturas versionadas; si se recapturan, se recalibran.
    """
    from PIL import Image as _Img
    with _Img.open(FIGS / nombre) as im:
        an, al = im.size
        iz, ar, de, ab = caja
        salida = FIGS / destino
        im.crop((int(an * iz), int(al * ar),
                 int(an * de), int(al * ab))).save(salida)
    return salida


def pulir_captura(ruta: Path):
    """
    Trata una captura como captura de producto: esquinas suaves y borde fino.

    Devuelve un PNG con canal alfa junto al original, con sufijo `_pulida`.
    Como el recorte lateral, se regenera en cada corrida y no se versiona
    (.gitignore): lo versionado es siempre la captura original.
    """
    from PIL import Image as _Img
    from PIL import ImageDraw as _Draw
    ruta = Path(ruta)
    with _Img.open(ruta) as im:
        im = im.convert("RGBA")
        an, al = im.size
        radio = max(10, an // 90)
        mascara = _Img.new("L", (an, al), 0)
        _Draw.Draw(mascara).rounded_rectangle((0, 0, an - 1, al - 1),
                                              radius=radio, fill=255)
        im.putalpha(mascara)
        _Draw.Draw(im).rounded_rectangle((0, 0, an - 1, al - 1), radius=radio,
                                         outline="#" + C["borde"], width=3)
        salida = FIGS / f"{ruta.stem}_pulida.png"
        im.save(salida)
    return salida


def imagen_encajada(s, ruta, x, y, w, h, centrar=True):
    """Mete la imagen en la caja (x,y,w,h) sin deformarla y la centra."""
    from PIL import Image as _Img
    iw, ih = _Img.open(ruta).size
    escala = min(w / iw, h / ih)
    aw, ah = iw * escala, ih * escala
    px = x + (w - aw) / 2 if centrar else x
    py = y + (h - ah) / 2 if centrar else y
    return s.shapes.add_picture(str(ruta), Inches(px), Inches(py),
                                width=Inches(aw), height=Inches(ah))


def rotulo(s, x, y, texto, color="acento_alto", tam=T_ETIQUETA):
    """Marca de anotación sobre una captura (lámina del formulario)."""
    fig = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                             Inches(0.06 + 0.098 * len(texto)), Inches(0.30))
    fig.adjustments[0] = 0.35
    fig.fill.solid()
    fig.fill.fore_color.rgb = rgb("acento_fondo")
    fig.line.color.rgb = rgb(color)
    fig.line.width = Pt(1)
    fig.shadow.inherit = False
    tf = fig.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    parrafo(tf, texto, tam=tam, color=color, negrita=True, fuente=MONO,
            primero=True, esp_despues=0, alin=PP_ALIGN.CENTER)
    return fig


# Las notas NO se escriben al construir: se registran junto con la fuente de
# la lámina y se vuelcan al final, dos veces — la versión ENTREGA lleva solo
# la línea de fuente, la versión EXPO lleva el guion completo del orador.
NOTAS_REG: list[tuple] = []
_ULTIMA_FUENTE = ""


def notas(s, texto, fuente=None):
    NOTAS_REG.append((s, texto,
                      _ULTIMA_FUENTE if fuente is None else fuente))


def pie_fuente(s, texto, y=None):
    global _ULTIMA_FUENTE
    _ULTIMA_FUENTE = texto
    y = Y_PIE if y is None else y
    tf = caja(s, MARGEN, y, UTIL, 0.34)
    parrafo(tf, "Fuente: " + texto, tam=T_PIE, color="texto_tenue",
            primero=True, esp_despues=0)


# ---------------------------------------------------------------------------
# Lámina 1 — carátula
# ---------------------------------------------------------------------------
APP_URL = "enaho-ingresos-informalidad.streamlit.app"
REPO_URL = "github.com/IchiSieben/enaho-ingresos-informalidad"

s = lamina()
TITULOS.append("Ingreso laboral e informalidad en el Perú (carátula)")
if (LOGOS / "inei_microdatos_cabecera.jpg").exists():
    imagen(s, LOGOS / "inei_microdatos_cabecera.jpg", MARGEN, 0.42, w=4.4)
if (LOGOS / "enei_logo.png").exists():
    imagen(s, LOGOS / "enei_logo.png", 11.05, 0.34, w=1.68)

tf = caja(s, MARGEN, 1.92, 12.09, 1.66)
parrafo(tf, "Ingreso laboral e informalidad en el Perú", tam=44, negrita=True,
        primero=True, esp_despues=6, esp_linea=0.94)
parrafo(tf, "Dos modelos sobre los microdatos de la ENAHO 2025 (INEI), "
            "desplegados en una app pública", tam=T_ENTRADILLA,
        color="texto_medio", esp_despues=0, esp_linea=1.05)

tf = caja(s, MARGEN, 3.62, 12.09, 0.62)
parrafo(tf, "Curso de Machine Learning · ENEI — Escuela Nacional de "
            "Estadística e Informática", tam=T_CUERPO, primero=True,
        esp_despues=3)
parrafo(tf, "Docente: Orlando Advíncula Zeballos · 25 de agosto de 2026",
        tam=T_CUERPO, color="texto_medio", esp_despues=0)

_, tf = panel(s, MARGEN, 4.48, 12.09, 1.16, relleno="acento_fondo",
              borde="acento")
parrafo(tf, "INTEGRANTES", tam=T_ETIQUETA, color="acento_alto", fuente=MONO,
        primero=True, esp_despues=5)
parrafo(tf, "Alan Nestor Cañazaca Mamani  ·  Magdalena Quico de la Cruz  ·  "
            "Yoichi Palacios Tanaka  ·  Edgar Delgado Ortega",
        tam=T_CUERPO, negrita=True, esp_despues=0, esp_linea=1.05)

# Dos líneas, alineadas a la izquierda: etiqueta en sans y URL en mono
# compartiendo línea base (por eso `parrafo_mixto` y no un párrafo por tramo).
# Son los dos únicos enlaces del mazo; ninguna otra lámina repite direcciones.
_, tf = panel(s, MARGEN, 5.86, 12.09, 1.24)
for _k, (_etq, _url) in enumerate((("App:  ", APP_URL),
                                   ("Repositorio:  ", REPO_URL))):
    # LEFT explícito: el primer párrafo de una autoshape hereda el centrado
    # del shape, y sin esto la línea de «App» salía corrida a la derecha
    # mientras la de «Repositorio» quedaba al margen.
    parrafo_mixto(tf, [(_etq, T_CUERPO, False, "texto_medio", FUENTE),
                       (_url, T_CUERPO, True, "acento_alto", MONO)],
                  primero=(_k == 0), alin=PP_ALIGN.LEFT,
                  esp_despues=10 if _k == 0 else 0)
# El hipervínculo va en una zona clicable INVISIBLE sobre cada URL, no en el
# run: PowerPoint restila todo run con hlinkClick (azul del tema y
# subrayado) aunque el rPr diga u="none" y traiga su propio color — se
# comprobó sobre el archivo generado. La acción de clic a nivel de shape
# abre el enlace igual en modo presentación y no toca el texto.
for _zx, _zy, _zw, _url in ((MARGEN + 0.50, 5.98, 6.00, APP_URL),
                            (MARGEN + 1.26, 6.42, 6.30, REPO_URL)):
    _z = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(_zx), Inches(_zy),
                            Inches(_zw), Inches(0.38))
    _z.fill.background()
    _z.line.fill.background()
    _z.shadow.inherit = False
    _z.click_action.hyperlink.address = "https://" + _url
notas(s, "QUÉ DIGO — Buenas, somos el grupo de Alan, Magdalena, Yoichi y "
         "Edgar. Hicimos dos modelos con los datos públicos de la encuesta "
         "de hogares del INEI: uno estima cuánto gana una persona al mes y "
         "otro qué tan probable es que su empleo sea informal. Los dos están "
         "funcionando ahora mismo en una página web pública que vamos a "
         "abrir en vivo. Como el curso es de despliegue, lo que más vamos a "
         "contar es cómo llega el modelo desde nuestra computadora hasta el "
         "navegador de cualquiera, y qué garantías tenemos de que lo que "
         "corre allá es lo mismo que entrenamos acá. Los dos enlaces de la "
         "pantalla son los únicos que hay que apuntar.\n\n"
         "TÉRMINOS DE ESTA LÁMINA — Microdatos: las respuestas individuales "
         "y anónimas de la encuesta, persona por persona, tal como las "
         "publica el INEI. App desplegada: la aplicación ya instalada en un "
         "servidor y accesible por una dirección web, no solo en nuestra "
         "máquina.\n\n"
         "SI PREGUNTAN — ¿La app está viva ahora mismo? Sí: es pública, "
         "corre en Streamlit Community Cloud y se abre desde el enlace de "
         "la carátula. En modo presentación el enlace es clicable.",
      fuente="AUTHORS.md · CITATION.cff · enlaces verificados en pantalla")

# ---------------------------------------------------------------------------
# Lámina 2 — qué construimos
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "Una app pública, dos modelos: cuánto gana un perfil y si su "
              "empleo es informal",
           "El mismo formulario alimenta dos modelos independientes; cada "
           "pestaña corre el suyo.")
# La base, en tres bloques que se leen de un vistazo (una franja corrida de
# texto no se leía): fuente · tamaño · variables. Cada cifra sale de su
# artefacto (embudo del informe, conteos del schema).
for _bx, _bw, _l1, _l2 in (
        (MARGEN, 3.70,
         [("ENAHO 2025 · INEI", 20, True, "acento_alto", FUENTE)],
         [("módulos ", T_CUERPO, False, "texto_medio", FUENTE),
          ("02 · 03 · 05", T_CUERPO, True, "texto", MONO)]),
        (MARGEN + 4.00, 4.09,
         [(EMBUDO["torneo"], 20, True, "acento_alto", MONO)],
         [("trabajadores con ingreso laboral", T_CUERPO, False,
           "texto_medio", FUENTE)]),
        (MARGEN + 8.39, 3.70,
         [(str(len(REG["features"])), 20, True, "acento_alto", MONO),
          (" variables", 20, True, "acento_alto", FUENTE)],
         [(f"{len(NUM)} numéricas · {len(CAT)} categóricas", T_CUERPO,
           False, "texto_medio", FUENTE)])):
    _, tf = panel(s, _bx, y, _bw, 0.92, relleno="acento_fondo",
                  borde="acento")
    parrafo_mixto(tf, _l1, primero=True, alin=PP_ALIGN.CENTER, esp_despues=1)
    parrafo_mixto(tf, _l2, alin=PP_ALIGN.CENTER, esp_despues=0)
ANCHO_CAP = 5.86
# Recortadas (la barra lateral no se lee proyectada) y pulidas como producto.
imagen_encajada(s, pulir_captura(recorte_lateral("cloud_reg_form.png",
                                                 "cloud_reg_panel.png")),
                MARGEN, y + 1.04, ANCHO_CAP, 2.30)
imagen_encajada(s, pulir_captura(recorte_lateral("cloud_clf_resultado.png",
                                                 "cloud_clf_panel.png")),
                6.85, y + 1.04, ANCHO_CAP, 2.30)
for x, etq, txt in (
        (MARGEN, "Estimación de ingreso  ·  regresión",
         "Estima el ingreso laboral mensual en soles con Gradient Boosting "
         "(especificación E9)."),
        (6.85, "Empleo informal  ·  clasificación",
         "Estima la probabilidad de que el empleo sea informal, con las "
         "mismas 11 variables.")):
    _, tf = panel(s, x, y + 3.46, ANCHO_CAP, 1.28)
    parrafo(tf, etq.upper(), tam=T_ETIQUETA, color="texto_tenue", fuente=MONO,
            primero=True, esp_despues=5)
    parrafo(tf, txt, tam=T_CUERPO, color="texto", esp_despues=0, esp_linea=1.04)
pie_fuente(s, f"INFORME_AUDITORIA.md §4 (embudo) · models/feature_schema.json "
              f"(variables) · capturas de la app desplegada ({APP_URL}).")
notas(s, f"QUÉ DIGO — Primero los datos, que son los tres bloques de "
         f"arriba: encuesta de hogares ENAHO 2025 del INEI, cruzando tres "
         f"de sus módulos — quiénes viven en el hogar, su educación y su "
         f"empleo. Después de los filtros quedan {EMBUDO['torneo']} "
         f"trabajadores con ingreso laboral. A cada uno lo describen once "
         f"variables, las mismas para los dos modelos; en la siguiente "
         f"lámina están una por una. Sobre esa base, la app responde dos "
         f"preguntas: a la izquierda, cuánto gana al mes un perfil como "
         f"este; a la derecha, qué tan probable es que ese empleo sea "
         f"informal. Son dos modelos separados que comparten las entradas, "
         f"no un modelo con dos salidas.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — Módulo: cada bloque temático de la "
         f"encuesta ({'02'} hogar, {'03'} educación, {'05'} empleo); se "
         f"cruzan por persona. Regresión: el modelo que responde con un "
         f"número (soles al mes). Clasificación: el modelo que responde con "
         f"una probabilidad entre 0 y 1 (informal o no). Modelo: el "
         f"programa que aprendió de los {EMBUDO['torneo']} casos a "
         f"responder esas preguntas para un perfil nuevo.\n\n"
         f"SI PREGUNTAN — ¿Por qué dos modelos y no uno? Porque son "
         f"preguntas de naturaleza distinta: una pide un monto y la otra "
         f"una probabilidad. Comparten las once entradas para que el mismo "
         f"formulario alimente a los dos, pero cada uno se entrenó y se "
         f"evalúa por su lado.")

# ---------------------------------------------------------------------------
# Lámina de variables (entrega: posición 3) — cuáles son y cómo se eligieron
# ---------------------------------------------------------------------------
s = lamina()
_F = {f["nombre"]: f for f in REG["features"]}
# Media línea por variable, armada desde el schema: etiqueta + rango (las
# numéricas) o niveles (las categóricas). Los rangos y conteos son datos del
# schema, así que el verificador los traza como cualquier otra cifra.
# Descripciones de UNA línea: una celda que se parte en dos hace crecer la
# tabla (python-pptx la CRECE, no la recorta) y el panel de abajo terminaba
# pisando las últimas filas.
_DESC = {
    "anios_educ": f"de {d(_F['anios_educ']['min'], 0)} a "
                  f"{d(_F['anios_educ']['max'], 0)} años aprobados",
    "edad": f"de {d(_F['edad']['min'], 0)} a {d(_F['edad']['max'], 0)} años",
    "exper": "edad − educación − 6",
    "exper2": "su cuadrado: efecto curvo",
    "horas_total": f"de {d(_F['horas_total']['min'], 0)} a "
                   f"{d(_F['horas_total']['max'], 0)} a la semana",
    "sexo": " / ".join(_F["sexo"]["opciones"]),
    "area": " / ".join(_F["area"]["opciones"]),
    "dominio": f"{len(_F['dominio']['opciones'])} regiones del país",
    "rama": f"{len(_F['rama']['opciones'])} sectores agrupados",
    "tamano_empresa": f"{len(_F['tamano_empresa']['opciones'])} tramos, "
                      f"desde «Hasta 20»",
    "categoria": f"{len(_F['categoria']['opciones'])} tipos de vínculo "
                 f"laboral",
}
_ETI_VAR = {"anios_educ": "Años de educación", "edad": "Edad",
            "exper": "Experiencia potencial", "exper2": "Experiencia²",
            "horas_total": "Horas semanales", "sexo": "Sexo",
            "area": "Área de residencia", "dominio": "Dominio geográfico",
            "rama": "Rama de actividad", "tamano_empresa": "Tamaño de empresa",
            "categoria": "Categoría ocupacional"}
y = titulo(s, "Las 11 variables: cuáles son y cómo quedaron elegidas",
           "Las mismas once entradas alimentan el regresor y el clasificador; "
           "un torneo de especificaciones decidió quedarse con ellas.")
filas_n = [[f"Numéricas ({len(NUM)})", "Qué es"]]
for f in NUM:
    filas_n.append([_ETI_VAR.get(f["nombre"], f["etiqueta"]),
                    _DESC[f["nombre"]]])
filas_c = [[f"Categóricas ({len(CAT)})", "Qué es"]]
for f in CAT:
    filas_c.append([_ETI_VAR.get(f["nombre"], f["etiqueta"]),
                    _DESC[f["nombre"]]])
tabla(s, MARGEN, y, 5.86, 1.95, filas_n, anchos=[44, 56], tam=16, tam_cab=16,
      cols_mono=[])
tabla(s, 6.85, y, 5.86, 2.28, filas_c, anchos=[44, 56], tam=16, tam_cab=16,
      cols_mono=[])
_, tf = panel(s, MARGEN, y + 2.90, 12.09, 1.56, relleno="acento_fondo",
              borde="acento")
parrafo(tf, "CÓMO QUEDARON ELEGIDAS", tam=T_ETIQUETA, color="acento_alto",
        fuente=MONO, primero=True, esp_despues=4)
parrafo(tf, f"Un torneo de nueve especificaciones — de la consigna del curso "
            f"(E1, S/ {d(TAB_TORNEO['E1']['MAE_cv'], 1)} de error al mes) a "
            f"variantes más ricas — compitió con los mismos datos y "
            f"pliegues. Ganó E9: Gradient Boosting sobre el logaritmo, con "
            f"estas once variables y S/ {d(TAB_TORNEO['E9']['MAE_cv'], 1)} "
            f"en validación cruzada. Nunca se eligió por el test.",
        tam=T_CUERPO, color="texto", esp_despues=0, esp_linea=1.05)
pie_fuente(s, "models/feature_schema.json (nombre, tipo, rango y niveles de "
              "cada variable) · models/ui_artifacts.json (torneo.tabla) · "
              "docs/METODOLOGIA_TORNEO.md")
notas(s, f"""QUÉ DIGO — Los dos modelos usan las mismas once variables: cinco numéricas y seis categóricas. Son cosas que cualquier persona sabe de sí misma: educación, edad, horas de trabajo, dónde vive, en qué sector trabaja y en qué tipo de empleo. No las elegimos a ojo: armamos un torneo de nueve recetas, desde la consigna del curso hasta variantes cada vez más ricas, y todas compitieron con los mismos datos y la misma regla. Ganó la novena, que usa estas once variables. El error bajó de S/ {d(TAB_TORNEO['E1']['MAE_cv'], 1)} al mes con la consigna a S/ {d(TAB_TORNEO['E9']['MAE_cv'], 1)} con la ganadora.

TÉRMINOS DE ESTA LÁMINA — Variable numérica: un número que entra tal cual (años, horas). Variable categórica: una opción de una lista cerrada (región, sector); el modelo la convierte por dentro en columnas de sí/no. Especificación: una receta concreta de modelo — qué variables entran y en qué forma. Torneo: todas las recetas compiten con los mismos datos y gana la que menos se equivoca. Validación cruzada: partir los datos en cinco, entrenar con cuatro partes y probar con la quinta, rotando, para medir sin hacerse trampa.

SI PREGUNTAN — ¿Por qué no está el tipo de contrato? Porque separa casi solo y se solapa con la propia definición de informalidad: el modelo aprendería la definición, no el fenómeno. Quedó fuera del clasificador y la decisión está documentada.""")

# ---------------------------------------------------------------------------
# Lámina 3 — la arquitectura, el ancla del mazo
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "Cada push a main redespliega la app en la nube, sin pasos "
              "manuales",
           "Del editor al navegador solo hay commits: Streamlit Cloud "
           "reconstruye y publica en cada push.")
# El diagrama manda: se lleva el grueso del alto y se centra en el ancho
# útil. Debajo, una banda de ancho completo con el mensaje de los
# microdatos — antes iba incrustado en un hueco del propio diagrama, donde
# flotaba suelto entre las cajas.
# La banda dice QUÉ sube y QUÉ no, sin ambigüedad: dos oraciones, ambas a
# 18 pt (el mínimo tipográfico del mazo no admite prosa menor) — la primera
# es la afirmación, la segunda el destino de lo que no sube.
ALTO_DIAG = ALTO_CUERPO - 1.54
imagen_encajada(s, F_ARQ, MARGEN, y, 12.09, ALTO_DIAG)
_, tf = panel(s, MARGEN, y + ALTO_DIAG + 0.14, 12.09, 1.32,
              relleno="buena_fondo", borde="buena")
parrafo(tf, f"Al repositorio sube el resumen, no la base: código, modelos ya "
            f"entrenados y resultados precalculados ({d(REPO['mb'], 2)} MB).",
        tam=T_CUERPO, color="buena", negrita=True, primero=True,
        esp_despues=3, esp_linea=1.04, alin=PP_ALIGN.CENTER)
parrafo(tf, f"Los {d(DATOS_MB, 1)} MB de microdatos se quedan en la máquina; "
            f"el README deja el enlace al portal del INEI para quien quiera "
            f"descargarlos y reproducir todo.",
        tam=T_CUERPO, color="texto", esp_despues=0, esp_linea=1.04,
        alin=PP_ALIGN.CENTER)
pie_fuente(s, "docs/arquitectura.md · .gitignore · git ls-files · tamaños "
              "medidos del disco al generar esta lámina")
notas(s, f"QUÉ DIGO — Esta lámina sostiene todo lo demás. Entrenamos en "
         f"nuestra computadora y de ahí sale un paquete chico: el código y "
         f"los resultados ya calculados, {d(REPO['mb'], 2)} megas en "
         f"{REPO['n_archivos']} archivos. Ese paquete se sube a GitHub con "
         f"un push. GitHub le avisa a Streamlit Cloud, y la nube "
         f"reconstruye la app y la publica sola. Entre guardar el cambio y "
         f"verlo en el navegador no hay ningún paso manual. Y los datos "
         f"originales del INEI, {d(DATOS_MB, 1)} megas, nunca salen de la "
         f"máquina: el README deja el enlace al portal del INEI solo para "
         f"quien quiera reproducir todo — la app nunca descarga nada del "
         f"INEI.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — Dónde vive cada cosa: TODO vive en "
         f"GitHub — el código, los modelos .joblib y los JSON viajan dentro "
         f"del repositorio; los microdatos no. Streamlit Cloud no guarda "
         f"nada propio: clona el repositorio, instala requirements.txt y "
         f"ejecuta. GitHub es la fuente; Streamlit solo lo corre y le pone "
         f"URL. Artefacto: un resultado ya calculado y guardado en un "
         f"archivo. .joblib: el modelo entrenado guardado como archivo. "
         f"Webhook: el aviso automático que GitHub le manda a Streamlit "
         f"cuando el repositorio cambia («hay versión nueva, redespliega»). "
         f"Repositorio: la carpeta versionada del proyecto, con su "
         f"historia completa.\n\n"
         f"SI PREGUNTAN — ¿Y si Streamlit Cloud se cae, o se borra hasta "
         f"la carpeta local? No se pierde nada: en Streamlit no vive nada "
         f"que no esté en GitHub — se reconecta el repositorio y la app "
         f"renace idéntica —, los microdatos se vuelven a bajar del portal "
         f"del INEI y los artefactos se regeneran corriendo el código. "
         f"Todo lo que no sube es recuperable o regenerable.")

# ---------------------------------------------------------------------------
# Lámina 4 — artefactos precomputados
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "La app no calcula al abrirse: lee resultados ya guardados",
           "Curvas, histogramas y tablas se calcularon una sola vez, al "
           "entrenar. La app solo las lee y las dibuja.")
DIAGRAMA_H = 3.10
# centrar=False: la figura se pega arriba de su caja (encajada por ancho
# mide ~2,89" de alto) y así la escalera de abajo respira sin robarle
# altura a las cajas grandes.
imagen_encajada(s, F_PRECOMP, MARGEN, y, 12.09, DIAGRAMA_H, centrar=False)
# La escalera de tamaños: tres órdenes de magnitud, de izquierda a derecha
# y con la cifra cada vez más chica — la base se queda, el resumen viaja,
# la app lee solo lo mínimo. Cada cifra sale de su medición.
Y_ESC = y + 3.04
for _px, _pw, _cifra, _tam_c, _col, _txt in (
        (MARGEN, 3.55, f"{d(DATOS_MB, 1)} MB", 20, "buena",
         "la base, se queda en casa"),
        (MARGEN + 4.04, 3.55, f"{d(REPO['mb'], 2)} MB", 19, "acento",
         "el resumen, viaja al repo"),
        (MARGEN + 8.08, 4.01, f"{d(UI_KB, 1)} KB", 18, "acento_alto",
         "lo que la app lee en cada visita")):
    _, tf = panel(s, _px, Y_ESC, _pw, 0.82)
    parrafo(tf, _cifra, tam=_tam_c, color=_col, negrita=True, fuente=MONO,
            primero=True, esp_despues=1, alin=PP_ALIGN.CENTER)
    parrafo(tf, _txt, tam=T_CUERPO, color="texto_medio", esp_despues=0,
            alin=PP_ALIGN.CENTER)
for _fx in (MARGEN + 3.55, MARGEN + 7.59):
    tf = caja(s, _fx, Y_ESC + 0.21, 0.49, 0.40)
    parrafo(tf, "→", tam=20, color="texto_tenue", negrita=True, fuente=MONO,
            primero=True, esp_despues=0, alin=PP_ALIGN.CENTER)
# El recorrido real, medido. UNA sola línea: el texto se acorta hasta que
# entra, nunca se parte dejando la cifra final huérfana. La banda termina
# a 0,16" del pie de fuente — el QA mide esa holgura.
_, tf_flujo = panel(s, MARGEN, Y_ESC + 1.00, 12.09, 0.56,
                    relleno="acento_fondo", borde="acento")
tf_flujo.vertical_anchor = MSO_ANCHOR.MIDDLE
T_FLUJO = 20
parrafo_mixto(tf_flujo, [
    ("Abrir la app", T_CUERPO, False, "texto_medio", FUENTE),
    ("  →  ", T_CUERPO, True, "texto_tenue", MONO),
    ("0", T_FLUJO, True, "buena", MONO),
    (" ejecuciones", T_CUERPO, False, "texto_medio", FUENTE),
    ("      Pulsar «Estimar»", T_CUERPO, False, "texto_medio", FUENTE),
    ("  →  ", T_CUERPO, True, "texto_tenue", MONO),
    ("1", T_FLUJO, True, "acento", MONO),
    (" predicción", T_CUERPO, False, "texto_medio", FUENTE),
    ("  →  ", T_CUERPO, True, "texto_tenue", MONO),
    (f"S/ {ING_TIPICO}", T_FLUJO, True, "acento_alto", MONO),
], primero=True, alin=PP_ALIGN.CENTER, esp_despues=0)
pie_fuente(s, "src/09_precomputar_ui.py · models/ui_artifacts.json · "
              "models/*.joblib · docs/presentacion/verificacion_local.py")
notas(s, f"QUÉ DIGO — Cuando alguien abre la app, el modelo no se ejecuta "
         f"ni una vez: todos los gráficos y tablas ya estaban calculados "
         f"desde el entrenamiento, guardados en un archivo de "
         f"{d(UI_KB, 1)} kilobytes que la app solo lee y dibuja. El modelo "
         f"trabaja únicamente cuando el usuario arma su perfil y pulsa "
         f"«Estimar»: una sola predicción, que para el perfil por defecto "
         f"da S/ {ING_TIPICO}. La escalera del medio es la decisión en "
         f"cifras: si subiéramos la base entera y calculáramos en el "
         f"servidor, la app cargaría lenta y se quedaría sin memoria; por "
         f"eso se precalcula una vez y en cada visita solo se lee.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — Precómputo: calcular una sola vez, al "
         f"entrenar, lo que la app va a mostrar siempre, en vez de "
         f"recalcularlo con cada visita. Artefacto: un resultado ya "
         f"calculado y guardado en un archivo. .joblib: el modelo entrenado "
         f"guardado como archivo; se carga y responde, sin reentrenar. "
         f"JSON: un archivo de texto con datos organizados, legible por "
         f"cualquier lenguaje. predict / predict_proba: la llamada que pide "
         f"al modelo un número (soles) o una probabilidad. Curva de umbral, "
         f"histograma, dependencia parcial: los gráficos de la app; los "
         f"tres vienen del precómputo. Caché: memoria de corto plazo de la "
         f"app para no releer archivos en cada clic.\n\n"
         f"SI PREGUNTAN — ¿Y si el perfil que arma el usuario no estaba "
         f"precalculado? Justo por eso los .joblib viajan con la app: el "
         f"perfil puntual sí se predice en vivo, es la única cuenta que se "
         f"hace en caliente y tarda milisegundos. Lo precalculado es todo "
         f"lo que no depende del perfil.")

# ---------------------------------------------------------------------------
# Lámina 5 — versiones fijadas
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "El pickle exige la misma scikit-learn que lo entrenó: cada "
              "versión queda fijada",
           "requirements.txt fija las versiones que generaron los artefactos, "
           "y Streamlit Cloud reconstruye ese mismo entorno.")
imagen_encajada(s, F_REQ, MARGEN, y, 8.30, 4.35)
tarjeta(s, 9.15, y + 0.08, 3.56, 2.08, "scikit-learn fijado", SKLEARN_REQ,
        "La misma que entrenó los .joblib. Si dejan de coincidir, este "
        "generador aborta.", color_cifra="acento")
tarjeta(s, 9.15, y + 2.26, 3.56, 2.08, "la nube apunta a", "main",
        "Streamlit Cloud sigue esa rama y reconstruye el entorno en "
        "cada push.", color_cifra="buena", tam_cifra=30)
pie_fuente(s, "requirements.txt · models/ui_artifacts.json (meta) · panel de "
              "Streamlit Community Cloud")
notas(s, f"QUÉ DIGO — Un modelo guardado solo se puede volver a abrir con "
         f"la misma versión de la librería que lo guardó. Si la nube "
         f"instalara otra, el modelo podría no cargar — o peor, cargar y "
         f"comportarse distinto sin avisar. Por eso requirements.txt fija "
         f"cada versión exacta: scikit-learn {SKLEARN_REQ}, la misma que "
         f"entrenó los modelos, y pandas {META['version_pandas']} por la "
         f"misma razón. Y la comprobación no es una promesa: el programa "
         f"que genera estas láminas compara ambos lados y se niega a "
         f"producir la presentación si no coinciden.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — Pickle / .joblib: el modelo entrenado "
         f"guardado como archivo binario; para abrirse necesita la misma "
         f"librería con la que se guardó. requirements.txt: la lista de "
         f"librerías con su versión exacta; la nube instala exactamente "
         f"eso. Entorno: el conjunto de librerías instaladas donde corre "
         f"la app; Streamlit Cloud lo reconstruye desde cero en cada "
         f"publicación. Artefacto: resultado ya calculado y guardado en un "
         f"archivo — aquí, los dos modelos y los JSON que la app lee.\n\n"
         f"SI PREGUNTAN — ¿Qué pasa si actualizan scikit-learn? Nada, "
         f"mientras requirements.txt no cambie: la nube instala la versión "
         f"fijada, no la última. Si algún día se quiere subir de versión, "
         f"hay que reentrenar, regenerar los artefactos y volver a "
         f"verificar — y el generador obliga a que todo eso pase junto.")

# ---------------------------------------------------------------------------
# Lámina 6 — el formulario por dentro
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "El formulario no declara ni un campo a mano: se genera desde "
              "feature_schema.json",
           "Cada componente de Streamlit resuelve un problema concreto; si el "
           "schema cambia, el formulario cambia solo.")
ANCHO_IMG = 6.90
imagen(s, pulir_captura(FIGS / "cloud_form.png"), MARGEN, y, w=ANCHO_IMG)
# Las chinchetas se colocan en fracciones de la captura, no en pulgadas
# sueltas: si mañana cambia el tamaño de la imagen, siguen en su sitio.
ALTO_IMG = ANCHO_IMG / 1.5
for k, (fx, fy) in enumerate(((0.10, 0.50), (0.30, 0.42), (0.30, 0.72),
                              (0.28, 0.22)), start=1):
    rotulo(s, MARGEN + fx * ANCHO_IMG, y + fy * ALTO_IMG, str(k))
# La leyenda va numerada y en una sola columna: cuatro entradas cortas caben,
# cuatro entradas explicadas no. El desarrollo largo está en las notas.
COMPONENTES = [
    ("1 · st.radio", "cambia de tema desde PALETAS"),
    ("2 · st.number_input", "una por variable numérica"),
    ("3 · st.selectbox", "una por variable categórica"),
    ("4 · st.expander · st.popover", "lectura llana y lectura técnica"),
]
tfc = caja(s, 7.75, y + 0.02, 4.96, 2.48)
for k, (nom, txt) in enumerate(COMPONENTES):
    parrafo(tfc, nom, tam=T_CUERPO, color="acento_alto", negrita=True,
            fuente=MONO, primero=(k == 0), esp_despues=1)
    parrafo(tfc, txt, tam=T_CUERPO, color="texto_medio", esp_despues=4,
            esp_linea=1.0)
# El control del umbral vive en la OTRA pestaña, así que no está en la
# captura de la izquierda: se capturó aparte de la app en producción
# (capturar_umbral.py) y va DENTRO de este recuadro, con su marcador sobre
# el slider real. El recuadro se dimensiona a partir del alto real de la
# imagen, no al tanteo: así nunca se sale por debajo del pie de fuente.
CAP_UMBRAL = FIGS / "cloud_umbral_control.png"
# Debajo de las cuatro entradas de la leyenda: cada una son dos renglones a
# 18 pt, así que la lista llega hasta ~y+2,6 y el recuadro tiene que
# empezar por debajo o le come la última línea al punto 4.
_Y_PANEL_5 = y + 2.80
if CAP_UMBRAL.exists():
    from PIL import Image as _ImgU
    _iw, _ih = _ImgU.open(CAP_UMBRAL).size
    _ancho_u = 4.62
    _alto_u = _ancho_u * _ih / _iw
    _x_u, _y_u = 7.92, _Y_PANEL_5 + 0.52
    panel(s, 7.75, _Y_PANEL_5, 4.96, _alto_u + 0.76,
          relleno="acento_fondo", borde="acento")
    tf = caja(s, 7.95, _Y_PANEL_5 + 0.14, 4.60, 0.32)
    parrafo(tf, "5 · st.slider + st.fragment", tam=T_CUERPO,
            color="acento_alto", negrita=True, fuente=MONO, primero=True,
            esp_despues=0)
    imagen(s, pulir_captura(CAP_UMBRAL), _x_u, _y_u, w=_ancho_u)
    # El marcador cae sobre el slider real, en fracciones de la captura:
    # el control ocupa la mitad derecha, a media altura.
    rotulo(s, _x_u + 0.66 * _ancho_u, _y_u + 0.44 * _alto_u, "5")
else:
    _, tf = panel(s, 7.75, _Y_PANEL_5, 4.96, 1.20, relleno="acento_fondo",
                  borde="acento")
    parrafo(tf, "5 · st.slider + st.fragment", tam=T_CUERPO,
            color="acento_alto", negrita=True, fuente=MONO, primero=True,
            esp_despues=2)
    parrafo(tf, "En «Empleo informal»: moverlo reejecuta solo ese fragmento.",
            tam=T_CUERPO, color="texto", esp_despues=0, esp_linea=1.02)
pie_fuente(s, "captura de la app desplegada · app/streamlit_app.py "
              "(componentes realmente usados) · models/feature_schema.json")
notas(s, "QUÉ DIGO — Este formulario no lo escribimos campo por campo: se "
         "fabrica solo. Hay un archivo que describe cada variable — cómo "
         "se llama, si es número u opción, su rango, sus niveles — y la "
         "app recorre ese mapa y crea un control por variable. Si mañana "
         "el modelo cambia una variable, el campo aparece o cambia solo, "
         "sin tocar la interfaz. Y es imposible que el formulario y el "
         "modelo se desincronicen, porque los dos nacen del mismo archivo. "
         "El control del umbral, en la otra pestaña, vive en un fragmento: "
         "moverlo redibuja solo ese bloque y no la página entera.\n\n"
         "TÉRMINOS DE ESTA LÁMINA — Schema (feature_schema.json): el JSON "
         "que describe cada variable — nombre, tipo, rango, niveles; la "
         "app recorre ese mapa y fabrica un control por variable; si el "
         "schema cambia, el campo aparece solo, y formulario y modelo no "
         "pueden desincronizarse porque nacen del mismo archivo. "
         "Componente: cada pieza de interfaz de Streamlit — número, "
         "selector, botón de opciones, deslizador. Fragment: un bloque de "
         "la página que se redibuja solo, sin recargar el resto. Expander "
         "y popover: los pliegues donde vive la lectura técnica sin "
         "estorbar la llana.\n\n"
         "SI PREGUNTAN — ¿Qué pasa si el usuario mete un valor fuera de "
         "rango? No puede: cada control nace con el rango y las opciones "
         "del schema, que son las que el modelo vio al entrenar. El "
         "formulario no deja escribir lo que el modelo no conoce.")

# ---------------------------------------------------------------------------
# Lámina 7 — disciplina de entrega
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "Lo que se publica lleva tests, cifras verificadas y medición "
              "sobre la app viva",
           "Tres controles rodean cada entrega: que la app y los archivos "
           "del modelo digan lo mismo, que cada cifra salga de su archivo, y "
           "qué se ve de verdad en la app.")
filas = [["Control", "Qué garantiza", "Dónde se comprueba"],
         ["Tests de contrato",
          "Que la app no pida nada que sus archivos no tengan. Los tres "
          "fallos reales están reproducidos sin navegador: si vuelven, "
          "saltan en local y no en la nube.",
          "tests/test_contratos.py"],
         ["Cifras verificadas",
          "Cada número de estas láminas se comprueba contra el archivo que "
          "lo genera; si una cifra no tiene fuente, la exposición no se "
          "genera.",
          "verificar_ppt.py"],
         ["QA sobre la app viva",
          "Lo que se afirma de la app desplegada se mide en su DOM con "
          "Playwright, no se mira a ojo.",
          "docs/POST_ENTREGA.md"]]
tabla(s, MARGEN, y, 12.09, 2.70, filas, anchos=[22, 52, 26], tam=16,
      tam_cab=16, cols_mono=[])
# Los tres fallos, nombrados. Un recuadro que dice «rompió tres contratos»
# sin decir cuáles no se puede narrar: cada línea dice qué se rompió y qué
# vio el usuario. Salen de tests/test_contratos.py, un test por fallo.
# Cada uno cabe en UNA línea a 18 pt: tres fallos de dos líneas desbordaban
# el panel y se metían debajo del pie de fuente.
_, tf = panel(s, MARGEN, y + 3.12, 12.09, 1.52, relleno="acento_fondo",
              borde="acento")
parrafo(tf, "LOS TRES FALLOS DEL 20 DE AGOSTO, HOY CADA UNO UN TEST",
        tam=T_ETIQUETA, color="acento_alto", fuente=MONO, primero=True,
        esp_despues=5)
for _fallo in (
        "Un tema del selector no estaba en la paleta: elegirlo dejaba la app "
        "en blanco, sin vuelta atrás.",
        "La app pedía un gráfico que el módulo desplegado no tenía: la "
        "página se cortaba a medio dibujo.",
        "Se renombró una clave del artefacto y la app siguió sirviendo la "
        "vieja desde su caché."):
    parrafo(tf, "—  " + _fallo, tam=T_CUERPO, color="texto", esp_despues=2,
            esp_linea=1.0)
pie_fuente(s, "tests/test_contratos.py · docs/presentacion/verificar_ppt.py "
              "· docs/POST_ENTREGA.md (estándares de la corrida de entrega)")
notas(s, "QUÉ DIGO — Tres controles, cada uno contra un tipo de error "
         "distinto. Primero, los tests de contrato: un despliegue del 20 "
         "de agosto rompió la app tres veces porque la app pedía cosas que "
         "sus archivos no tenían; hoy cada uno de esos fallos es una "
         "prueba automática que salta en nuestra máquina antes de "
         "publicar. Segundo, las cifras: cada número de estas láminas se "
         "comprueba contra el archivo que lo genera, y si un número no "
         "tiene fuente, la presentación no se genera. Tercero, la app "
         "viva: lo que afirmamos de ella no se mira a ojo, se mide sobre "
         "la página real con un navegador automatizado.\n\n"
         "TÉRMINOS DE ESTA LÁMINA — Test de contrato: una prueba de que "
         "dos piezas siguen encajando — que la app no pida nada que sus "
         "archivos no tengan. Caché: la memoria de corto plazo de la app; "
         "uno de los fallos fue servir un archivo viejo desde ahí. Cifras "
         "verificadas: cada número visible (y el de estas notas) se traza "
         "hasta el archivo que lo produce con un script auditor. DOM: la "
         "estructura interna de la página web ya dibujada; medir ahí es "
         "medir lo que el usuario ve de verdad. Playwright: la herramienta "
         "que abre la app real y la inspecciona sola.\n\n"
         "SI PREGUNTAN — ¿Quién verifica las cifras de ESTA exposición? El "
         "mismo estándar: la presentación se genera desde los artefactos y "
         "un script la reabre y comprueba número por número, incluidas "
         "estas notas. Lo que están viendo pasó esa auditoría.")

# ---------------------------------------------------------------------------
# Lámina 8 — los modelos: RF y GB en los dos problemas (guía)
# ---------------------------------------------------------------------------
s = lamina()
e7, e8, e9 = TAB_TORNEO["E7"], TAB_TORNEO["E8"], TAB_TORNEO["E9"]
gb = next(f for f in CLF_COMP if "Gradient" in f["algoritmo"])
rf = next(f for f in CLF_COMP if "Random" in f["algoritmo"])
lo = next(f for f in CLF_COMP if f not in (gb, rf))
y = titulo(s, "Elegimos Gradient Boosting: le gana a Random Forest por poco, "
              "con las mismas 11 variables",
           "Mismos datos, mismos cinco pliegues, misma semilla; el modelo se "
           "eligió por validación cruzada, nunca por el test.")
filas = [["", "Random Forest", "Gradient Boosting", "Referencia lineal"],
         ["Ingreso mensual — MAE cv, soles",
          d(e8["MAE_cv"], 1), d(e9["MAE_cv"], 1),
          f"{d(e7['MAE_cv'], 1)}  (E7)"],
         ["Empleo informal — PR-AUC cv",
          d(rf["PRAUC_cv"], 4), d(gb["PRAUC_cv"], 4),
          f"{d(lo['PRAUC_cv'], 4)}  (logística)"]]
tabla(s, MARGEN, y, 12.09, 1.60, filas, anchos=[34, 21, 24, 21],
      resaltar_col=2, cols_mono=[1, 2, 3])
filas_v = [["Variables de entrada",
            f"las mismas {len(NUM) + len(CAT)} variables en los dos modelos"],
           [f"Numéricas ({len(NUM)})",
            "años de educación · edad · experiencia · experiencia² · "
            "horas semanales"],
           [f"Categóricas ({len(CAT)})",
            f"sexo · área · dominio geográfico · rama de actividad · tamaño "
            f"de empresa · categoría ocupacional — {N_NIVELES} niveles "
            f"one-hot dentro del pipeline"]]
tabla(s, MARGEN, y + 2.02, 12.09, 1.50, filas_v, anchos=[24, 76], tam=16,
      tam_cab=16, cols_mono=[])
_, tf = panel(s, MARGEN, y + 3.74, 12.09, 0.86, relleno="acento_fondo",
              borde="acento")
parrafo(tf, f"El clasificador no se corta en 0,5: el umbral operativo "
            f"{d(PUNTO['umbral'], 4)} se fijó para que la precisión llegue a "
            f"{d(PUNTO['precision_oof'], 2)} — esa es una decisión de costos, "
            f"no del modelo.",
        tam=T_CUERPO, color="texto", primero=True, esp_despues=0,
        esp_linea=1.06)
pie_fuente(s, "models/ui_artifacts.json (torneo.tabla, "
              "clasificador.comparacion) · models/feature_schema.json · "
              "src/04, src/06")
notas(s, f"QUÉ DIGO — Probamos las dos familias de árboles del curso en "
         f"los dos problemas, con los mismos datos y las mismas "
         f"particiones. En ingreso, Gradient Boosting se equivoca en "
         f"promedio S/ {d(e9['MAE_cv'], 1)} al mes y Random Forest "
         f"S/ {d(e8['MAE_cv'], 1)}: un empate práctico, con el mejor "
         f"modelo lineal lejos, en S/ {d(e7['MAE_cv'], 1)}. En "
         f"informalidad también gana Gradient Boosting por poco. Elegimos "
         f"siempre con validación cruzada; el conjunto de prueba se miró "
         f"una sola vez, al final. Y el corte del clasificador no es el "
         f"0,5 por defecto: lo pusimos donde la precisión llega a "
         f"{d(PUNTO['precision_oof'], 2)}, que es una decisión de costos, "
         f"no del modelo.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — MAE: en cuánto se equivoca el modelo "
         f"en promedio, en soles al mes. PR-AUC: qué tan bien separa "
         f"informal de formal, de 0 a 1; la usamos en vez del acierto "
         f"simple porque el {d(BASE_PR * 100, 1)} % ya es informal y "
         f"acertar diciendo siempre «informal» sería trampa. Validación "
         f"cruzada: partir los datos en cinco, entrenar con cuatro y "
         f"probar con el quinto, rotando; mide sin hacerse trampa. Random "
         f"Forest / Gradient Boosting: dos maneras de combinar muchos "
         f"árboles de decisión — en paralelo la primera, corrigiéndose en "
         f"cadena la segunda. Umbral: la probabilidad mínima a partir de "
         f"la cual el clasificador señala un empleo. Semilla: el número "
         f"fijo que hace reproducible el azar de las particiones.\n\n"
         f"SI PREGUNTAN — ¿De dónde salen los {EMBUDO['torneo']} casos? "
         f"Del embudo documentado: {EMBUDO['crudo']} registros del módulo "
         f"de empleo, {EMBUDO['ocupados']} ocupados de 14 años o más, "
         f"{EMBUDO['modelado']} con ingreso positivo y {EMBUDO['torneo']} "
         f"completos, partidos en {EMBUDO['train']} de entrenamiento y "
         f"{EMBUDO['test']} de prueba con semilla fija.")

# ---------------------------------------------------------------------------
# Lámina 9 — qué variables pesan (guía)
# ---------------------------------------------------------------------------
s = lamina()
_orden = sorted(zip(IMP_REG["variables"], IMP_REG["media"]), key=lambda t: -t[1])
_orden_clf = sorted(zip(IMP_CLF["variables"], IMP_CLF["media"]),
                    key=lambda t: -t[1])
# El título afirma que el trío es el mismo en ambos modelos. Se comprueba,
# no se supone: si algún día los artefactos divergen, no se genera una
# lámina que afirma algo que dejó de ser cierto.
if ({v for v, _ in _orden[:3]} != {v for v, _ in _orden_clf[:3]}):
    raise SystemExit(
        "El top 3 de importancia ya no coincide entre regresor y "
        "clasificador: el título de la lámina de variables lo afirma. "
        "Revísalo antes de generar.")
y = titulo(s, "Las mismas tres variables cargan los dos modelos: categoría, "
              "educación y tamaño de empresa",
           "Con cada modelo ya entrenado, desordenamos una columna a la vez "
           "y medimos cuánto empeora: soles de error en el regresor, puntos "
           "de PR-AUC en el clasificador. Uso, no causa.")
imagen_encajada(s, F_IMP, MARGEN, y, 8.00, 4.55)
tarjeta(s, 8.86, y + 0.02, 3.85, 1.28, "categoría ocupacional",
        f"S/ {n(_orden[0][1])}", "Soles de error si se pierde.",
        color_cifra="acento", tam_cifra=28)
tarjeta(s, 8.86, y + 1.42, 3.85, 1.28, "años de educación",
        f"S/ {n(_orden[1][1])}", "La segunda, mismo cálculo.",
        color_cifra="acento", tam_cifra=28)
# El clasificador, con su propia métrica en el propio encabezado del panel.
# El valor va PRIMERO y en mono: así los tres quedan alineados en columna
# en vez de dejar un borde derecho irregular.
_ETI_CORTA = {"tamano_empresa": "tamaño de empresa",
              "categoria": "categoría ocupacional",
              "anios_educ": "años de educación",
              "horas_total": "horas semanales"}
# El encabezado integra la métrica en DOS líneas exactas de mono 14
# (~30 caracteres por línea en 3,45" útiles): más largo pasaba a tres
# líneas y empujaba la última fila fuera del panel, sobre el pie.
_, tf = panel(s, 8.86, y + 2.82, 3.85, 1.76, relleno="acento_fondo",
              borde="acento")
parrafo(tf, "CLASIFICADOR · TOP 3 — CAÍDA DE PR-AUC AL DESORDENAR",
        tam=T_ETIQUETA, color="acento_alto", fuente=MONO,
        primero=True, esp_despues=6, esp_linea=1.05)
for _nom, _med in _orden_clf[:3]:
    parrafo_mixto(tf, [
        (f"−{d(_med, 3)}", T_CUERPO, True, "acento_alto", MONO),
        (f"  {_ETI_CORTA.get(_nom, _nom)}", T_CUERPO, False, "texto", FUENTE),
    ], esp_despues=4, esp_linea=1.0)
pie_fuente(s, f"models/ui_artifacts.json (importancia_permutacion del "
              f"regresor y del clasificador): {IMP_REG['n_repeticiones']} "
              f"repeticiones sobre {n(IMP_REG['n_filas'])} filas de test en "
              f"ambos")
notas(s, f"QUÉ DIGO — Con cada modelo ya entrenado hicimos el mismo "
         f"experimento: desordenar una columna al azar y medir cuánto "
         f"empeora. En el regresor se mide en soles: sin la categoría "
         f"ocupacional el error sube unos S/ {n(_orden[0][1])} al mes, y "
         f"sin la educación unos S/ {n(_orden[1][1])}. En el clasificador "
         f"se mide en puntos de PR-AUC, y salen las mismas tres variables "
         f"en otro orden. Importante: esto dice cuánto USA el modelo cada "
         f"variable, no qué causa el ingreso o la informalidad en la vida "
         f"real.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — Permutación: desordenar una columna "
         f"para ver cuánto empeora el modelo sin ella; si no la usaba, no "
         f"pasa nada. Puntos de PR-AUC: el clasificador acierta con PR-AUC "
         f"de {d(gb['PRAUC_cv'], 2)}; al desordenar «tamaño de empresa» "
         f"cae {d(_orden_clf[0][1], 3)} — es el mismo experimento que los "
         f"soles, con otra regla de medir. Y es el mismo trío de variables "
         f"en los dos modelos, en distinto orden. PR-AUC: qué tan bien "
         f"separa informal de formal, de 0 a 1. Uso vs. causa: que el "
         f"modelo se apoye en una variable no significa que esa variable "
         f"cause el resultado; la lectura causal vive en el modelo "
         f"explicativo de la app.\n\n"
         f"SI PREGUNTAN — ¿Esto coincide con la ablación que hicieron? Sí: "
         f"reentrenar sin tamaño de empresa ni categoría baja el PR-AUC de "
         f"{d(float(ABLACION[0]['PRAUC_cv']), 4)} a "
         f"{d(float(ABL_V2['PRAUC_cv']), 4)} — dos métodos distintos, "
         f"misma conclusión: esas variables llevan la señal, coherente con "
         f"cómo se define la informalidad.")

# ---------------------------------------------------------------------------
# Lámina 10 — verificación local = nube, los dos modelos (guía)
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "El modelo local y el desplegado dan el mismo número, dígito a "
              "dígito",
           "El mismo perfil se corrió en la consola de VS Code y en la app "
           "pública: tres números, cero diferencias.")
tf = caja(s, MARGEN, y + 0.02, 5.30, 0.26)
parrafo(tf, "VS CODE · CONSOLA LOCAL", tam=T_ETIQUETA, color="texto_tenue",
        fuente=MONO, primero=True, esp_despues=0, alin=PP_ALIGN.CENTER)
tf = caja(s, 6.62, y + 0.02, 6.09, 0.26)
parrafo(tf, "STREAMLIT COMMUNITY CLOUD · APP PÚBLICA", tam=T_ETIQUETA,
        color="texto_tenue", fuente=MONO, primero=True, esp_despues=0,
        alin=PP_ALIGN.CENTER)
# De cada captura de la app se proyecta SOLO la zona del número comparado:
# las tarjetas del regresor y la banda del veredicto del clasificador.
ZONA_REG = pulir_captura(recorte_caja(
    "cloud_reg_tarjetas.png", "cloud_reg_zona.png",
    (0.525, 0.440, 0.955, 0.760)))
ZONA_CLF = pulir_captura(recorte_caja(
    "cloud_clf_resultado.png", "cloud_clf_zona.png",
    (0.525, 0.010, 0.968, 0.198)))
for fila_y, fig_con, captura in ((y + 0.34, F_CONSOLA_REG, ZONA_REG),
                                 (y + 1.86, F_CONSOLA_CLF, ZONA_CLF)):
    imagen_encajada(s, fig_con, MARGEN, fila_y, 5.30, 1.42)
    tf = caja(s, 5.98, fila_y + 0.46, 0.60, 0.50)
    parrafo(tf, "=", tam=40, color="buena", negrita=True, primero=True,
            esp_despues=0, alin=PP_ALIGN.CENTER)
    imagen_encajada(s, captura, 6.62, fila_y, 6.09, 1.42)
TARJETAS = [
    ("ingreso típico (mediana)", f"S/ {ING_TIPICO}",
     "Idéntico en los dos lados.", "acento"),
    ("ingreso esperado (media)", f"S/ {ING_ESPERADO}",
     f"La mediana × Duan {d(REG['smearing_duan'], 4)}.", "acento"),
    ("prob. de informalidad", PROBA_ES,
     f"Umbral {d(PUNTO['umbral'], 4)}: señalado.", "media"),
]
_anch = (12.09 - 0.35 * (len(TARJETAS) - 1)) / len(TARJETAS)
for k, (etq, cif, txt, col) in enumerate(TARJETAS):
    tarjeta(s, MARGEN + k * (_anch + 0.35), y + 3.50, _anch, 1.22,
            etq, cif, txt, color_cifra=col, tam_cifra=22)
pie_fuente(s, "docs/presentacion/verificacion_local.py · "
              "salida_consola_verificacion.txt · capturas de la app desplegada")
_horas_reg = next(f["default"] for f in REG["features"]
                  if f["nombre"] == "horas_total")
_horas_clf = next(f["default"] for f in CLF["features"]
                  if f["nombre"] == "horas_total")
notas(s, f"QUÉ DIGO — Antes de dar la app por buena corrimos el mismo "
         f"perfil en dos lugares: la consola de nuestra computadora y la "
         f"app ya publicada. Los números coinciden dígito a dígito: "
         f"S/ {ING_TIPICO} de ingreso típico, S/ {ING_ESPERADO} de ingreso "
         f"esperado y {PROBA_ES} de probabilidad de informalidad. Ojo con "
         f"lo que prueba esto: no dice que el modelo prediga bien — eso ya "
         f"se midió —, dice que la nube corre EXACTAMENTE el modelo que "
         f"entrenamos, no una copia vieja ni un reentrenamiento distinto.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — Mediana: el valor del medio — la "
         f"mitad de los perfiles como este gana menos de S/ {ING_TIPICO}, "
         f"la otra mitad más. Factor de Duan: la corrección "
         f"({d(REG['smearing_duan'], 4)}) por haber entrenado el modelo en "
         f"logaritmo del ingreso; convierte la mediana en el promedio "
         f"esperado, que es más alto porque unos pocos sueldos grandes lo "
         f"jalan. Umbral: la vara de decisión; con {d(PUNTO['umbral'], 4)} "
         f"este perfil queda señalado. Consola: la salida de texto del "
         f"script en la computadora, sin interfaz. Artefacto: el modelo "
         f"guardado que viaja al repositorio — el mismo archivo en los dos "
         f"lados, por eso los números no pueden diferir.\n\n"
         f"SI PREGUNTAN — ¿Por qué la consola dice INFORMAL y la app "
         f"«señalado para focalización»? Es lo mismo dicho con cuidado: la "
         f"señal es sobre una configuración de empleo, no un juicio sobre "
         f"la persona. Y los dos formularios arrancan con horas distintas "
         f"({d(_horas_reg, 0)} y {d(_horas_clf, 0)}), así que las dos "
         f"filas no comparan entre sí — cada fila compara consola contra "
         f"app con el mismo perfil.")

# ---------------------------------------------------------------------------
# Lámina 11 — la auditoría
# ---------------------------------------------------------------------------
s = lamina()
_aut = UA["torneo"]["autopsia"]
y = titulo(s, "Revisamos nuestro propio trabajo: cuatro errores encontrados, "
              "cuatro publicados",
           "Cada cifra contra el archivo que la genera, cada cita contra su "
           "fuente; cada hallazgo, con su origen.")
# La columna «Qué pasó» se narra: cada fila dice qué se leía mal y en qué
# quedó. El texto sigue a la sección «Qué encontró la auditoría» de la app
# (app/streamlit_app.py · AUDITORIA), que es donde vive la versión canónica.
filas = [["Hallazgo", "Dónde nació", "Qué pasó"],
         ["El código de faltante, leído como un sueldo",
          "DATOS DE ORIGEN",
          f"999999 es «no sabe», pero entraba como S/ 999.999: tratarlo como "
          f"faltante subió el R² de {d(_aut['corrida_sucia']['r2'], 3)} a "
          f"{d(_aut['corrida_limpia']['r2'], 3)}"],
         ["La rejilla de hiperparámetros estaba acotada",
          "DECISIÓN PROPIA",
          f"Los tres quedaron en el borde. Ampliada, el error baja de "
          f"S/ {REJILLA['vieja']} a S/ {REJILLA['nueva']} "
          f"({REJILLA['mejora_pct']} %): no se promovió"],
         ["Una cifra del INEI con la etiqueta equivocada",
          "DECISIÓN PROPIA",
          f"El {EXTERNA['inei_1_10']} % es del tramo 1-10 del INEI; nuestra "
          f"categoría «Hasta 20» da {d(TASA_HASTA20, 1)} %"],
         ["Tres afirmaciones distintas sobre el mismo R²",
          "DECISIÓN PROPIA · DOCUMENTACIÓN",
          "Circulaban «0,4–0,5», «rara vez supera 0,4» y «ninguno supera "
          "0,5»: quedó una, sobre Mincer y Card"]]
# Ojo: python-pptx NO recorta una tabla, la CRECE hasta que quepa el texto. El
# alto que se le pasa es un mínimo, así que lo que va debajo hay que colocarlo
# contando con ese crecimiento — o acortar la celda, que es lo que se hizo.
tabla(s, MARGEN, y, 12.09, 2.65, filas, anchos=[30, 23, 47], tam=16,
      tam_cab=16, cols_mono=[])
# La tabla CRECE con el texto: con «Qué pasó» narrado, la cuarta fila baja
# hasta ~y+3,3. El panel arranca por debajo de eso o le pisa la última fila.
_, tf = panel(s, MARGEN, y + 3.44, 12.09, 1.20, relleno="acento_fondo",
              borde="acento")
parrafo(tf, "LA LÍNEA DE MÉTODO", tam=T_ETIQUETA, color="acento_alto",
        fuente=MONO, primero=True, esp_despues=4)
parrafo(tf, "Los de origen se corrigen y se documentan; los propios se "
            "corrigen y se aprende de ellos; los de cita se verifican en el "
            "texto completo. Ninguno se borra en silencio.",
        tam=T_CUERPO, color="texto", esp_despues=0, esp_linea=1.04)
pie_fuente(s, "INFORME_AUDITORIA.md · models/ui_artifacts.json (torneo.autopsia) "
              "· la sección «Qué encontró la auditoría» de la propia app")
notas(s, f"QUÉ DIGO — Antes de publicar revisamos nuestro propio trabajo: "
         f"cada cifra contra el archivo que la genera, cada cita contra su "
         f"fuente. Encontramos cuatro errores y publicamos los cuatro, con "
         f"la etiqueta de dónde nació cada uno. El más grave venía en los "
         f"datos: el código «no sabe» del INEI entraba como un sueldo de "
         f"casi un millón de soles y torcía toda la regresión. Los otros "
         f"tres los pusimos nosotros al elegir o al resumir. Ninguno cambia "
         f"el modelo publicado, y ninguno se borró en silencio: el valor de "
         f"auditar no fue no tener errores, fue encontrarlos y dejarlos "
         f"escritos. El detalle narrado de cada uno está en "
         f"guion_hallazgos.md.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — Centinela: un código que la encuesta "
         f"usa para «no sabe» (999999); si se lee como número real, "
         f"envenena el modelo. R²: qué parte de la variación del ingreso "
         f"explica el modelo, de 0 a 1 — aquí saltó de "
         f"{d(_aut['corrida_sucia']['r2'], 3)} a "
         f"{d(_aut['corrida_limpia']['r2'], 3)} al limpiar el centinela. "
         f"Rejilla de hiperparámetros: la lista de configuraciones que se "
         f"prueban al ajustar el modelo; si el ganador queda en el borde, "
         f"el óptimo puede estar fuera. Tramo: el rango de tamaño de "
         f"empresa con el que el INEI reporta ({EXTERNA['inei_1_10']} % es "
         f"del tramo 1-10, no de nuestra categoría «Hasta 20», que da "
         f"{d(TASA_HASTA20, 1)} %).\n\n"
         f"SI PREGUNTAN — ¿Por qué no promovieron la rejilla mejor si gana "
         f"en los cinco pliegues? Porque la mejora es de "
         f"{REJILLA['mejora_pct']} % — real pero irrelevante en la "
         f"práctica: no justifica regenerar el modelo publicado, revalidar "
         f"la corrección de Duan y rehacer el precómputo. Separar «es "
         f"mejor» de «vale la pena cambiarlo» es parte del método.")

# ---------------------------------------------------------------------------
# Lámina de prácticas (entrega: posición 13) — el inventario, con su dónde
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "Las prácticas que sostienen el código",
           "Seis prácticas de ingeniería, cada una con su lugar concreto en "
           "este proyecto y el error que evita.")
imagen_encajada(s, F_FUENTE, MARGEN, y, 12.09, 1.30)
# Celdas de UNA línea: las tablas de python-pptx crecen con el texto y ya
# nos pisaron un panel una vez. Los anchos están medidos para eso.
filas = [["Práctica", "Dónde vive en este proyecto", "Qué error evita"],
         ["Fuente única (DRY)",
          "schema→formulario · artefactos→PPT · requirements→nube",
          "Dos copias que se contradicen"],
         ["Precomputar lo pesado",
          f"src/09 → ui_artifacts.json ({d(UI_KB, 1)} KB), la app solo lee",
          "App lenta o sin memoria"],
         ["Versiones fijadas",
          "requirements.txt = meta del artefacto, o se aborta",
          "Un modelo distinto en la nube"],
         ["Cada bug, un test",
          "test_contratos.py reproduce los tres fallos del 20/08",
          "Que el mismo error regrese"],
         ["Despliegue sin manos",
          "push → webhook → redeploy automático",
          "Pasos manuales que se olvidan"],
         ["Verificar, no confiar",
          "local = desplegado · verificar_ppt.py traza cada cifra",
          "Números dichos de memoria"]]
tabla(s, MARGEN, y + 1.48, 12.09, 2.55, filas, anchos=[20, 50, 30], tam=16,
      tam_cab=16, cols_mono=[])
pie_fuente(s, "models/feature_schema.json · models/ui_artifacts.json · "
              "requirements.txt · tests/test_contratos.py · "
              "docs/presentacion/verificar_ppt.py")
notas(s, f"""QUÉ DIGO — Nada de esto lo inventamos: son prácticas estándar de ingeniería, y esta lámina dice dónde vive cada una en el proyecto. Una sola fuente de verdad: el formulario, estas láminas y la nube leen archivos generados; nadie copia números a mano. Lo pesado se precalcula una vez. Las versiones van fijadas y el generador aborta si divergen. Cada bug que tuvimos quedó convertido en un test. El despliegue no tiene pasos manuales. Y no confiamos: verificamos — el modelo local contra el desplegado, y cada cifra contra su archivo. La corriente va en una sola dirección: entrenamiento → schema → formulario; nadie edita la copia, todos leen la fuente.

TÉRMINOS DE ESTA LÁMINA — DRY («don't repeat yourself»): cada dato vive en UN solo lugar y los demás lo leen de ahí — es la caja del diagrama. Fallo silencioso: el error que no revienta; si el formulario escribiera «hasta 20 trabajadores» donde el modelo aprendió «Hasta 20», el modelo no explota — predice basura sin avisar. Con el schema ese error no puede existir, porque el menú y el modelo leen el mismo archivo. Antipatrón: la solución que se sabe que sale mal — aquí, mantener dos copias del mismo dato (el lado tachado del diagrama). Webhook: el aviso automático de GitHub a la nube.

SI PREGUNTAN — ¿Y si alguien quiere cambiar una opción del formulario? No puede hacerlo a mano: tendría que cambiar el schema, que sale del entrenamiento — y entonces el formulario, la app y estas láminas cambian juntos. Esa es exactamente la gracia de la fuente única.""")

# ---------------------------------------------------------------------------
# Lámina 12 — cierre
# ---------------------------------------------------------------------------
s = lamina()
y = titulo(s, "El regresor compara perfiles; el clasificador señala empleos, "
              "no personas",
           "Qué se puede afirmar con cada modelo, qué no, y qué dejó el "
           "despliegue.")
filas = [["", "Regresor de ingreso", "Clasificador de informalidad"],
         ["Sí se puede afirmar",
          f"Un ingreso típico para un perfil, con un error medio de "
          f"S/ {d(REG['metricas_test']['mae_mediana'], 0)} al mes.",
          f"Que una configuración laboral observable se asocia a informalidad "
          f"(PR-AUC de test {d(CLF['metricas_test']['prauc'], 4)})."],
         ["No se puede afirmar",
          "Cuánto va a cobrar una persona concreta. Sirve para comparar "
          "perfiles, no para fijar sueldos.",
          "Que alguien vaya a ser informal. No predice el futuro ni juzga a "
          "ninguna persona."]]
tabla(s, MARGEN, y, 12.09, 2.10, filas, anchos=[20, 40, 40], tam=16,
      tam_cab=16, cols_mono=[])
_, tf = panel(s, MARGEN, y + 2.50, 12.09, 1.26, relleno="acento_fondo",
              borde="acento")
parrafo(tf, "QUÉ DEJÓ EL DESPLIEGUE", tam=T_ETIQUETA, color="acento_alto",
        fuente=MONO, primero=True, esp_despues=6)
parrafo(tf, f"Fijar las versiones exactas (scikit-learn {SKLEARN_REQ}, pandas "
            f"{META['version_pandas']}) y precomputar todo lo pesado en un "
            f"JSON de {d(UI_KB, 1)} KB es lo que hace que la app cargue "
            f"rápido y no se quede sin memoria.",
        tam=T_CUERPO, color="texto", esp_despues=0, esp_linea=1.06)
tf = caja(s, MARGEN, y + 4.02, 12.09, 0.62)
parrafo(tf, "Alan Nestor Cañazaca Mamani · Magdalena Quico de la Cruz · "
            "Yoichi Palacios Tanaka · Edgar Delgado Ortega",
        tam=T_CUERPO, color="texto", negrita=True, primero=True,
        esp_despues=2, alin=PP_ALIGN.CENTER)
parrafo(tf, "Docente: Orlando Advíncula Zeballos",
        tam=T_CUERPO, color="texto_medio", esp_despues=0, alin=PP_ALIGN.CENTER)
pie_fuente(s, "models/feature_schema.json (métricas de test) · AUTHORS.md · "
              "CITATION.cff · LICENSE (Apache-2.0) y docs/LICENSE-DOCS.md")
notas(s, f"QUÉ DIGO — Para cerrar, los límites, que también son parte del "
         f"trabajo. El regresor compara perfiles: dice el ingreso típico de "
         f"gente como la del perfil, con un error promedio de unos "
         f"S/ {d(REG['metricas_test']['mae_mediana'], 0)} al mes; no le "
         f"dice a nadie cuánto va a cobrar. El clasificador señala "
         f"configuraciones de empleo que se asocian a informalidad; no "
         f"predice el futuro ni juzga personas. Del despliegue nos llevamos "
         f"dos aprendizajes: fijar las versiones exactas, porque un modelo "
         f"guardado no viaja entre versiones, y dejar todo lo pesado "
         f"calculado de antemano. Los enlaces están en la portada. "
         f"Gracias.\n\n"
         f"TÉRMINOS DE ESTA LÁMINA — Error medio (MAE): en cuánto se "
         f"equivoca el modelo en promedio, en soles al mes. PR-AUC de "
         f"test ({d(CLF['metricas_test']['prauc'], 4)}): qué tan bien "
         f"separa informal de formal, de 0 a 1, medido en datos que el "
         f"modelo nunca vio. Focalización: usar la señal para priorizar a "
         f"quién mirar primero, no para etiquetar personas. Precomputar: "
         f"dejar calculado de antemano lo que la app va a mostrar "
         f"siempre.\n\n"
         f"SI PREGUNTAN — ¿Esto se podría usar en producción de verdad? "
         f"Como herramienta demostrativa y de priorización, sí — es "
         f"pública y reproducible. Para decisiones sobre personas "
         f"concretas, no: el propio mazo dice qué se puede afirmar y qué "
         f"no, y esa línea no se cruza.")

# ---------------------------------------------------------------------------
# Guardado
# ---------------------------------------------------------------------------
GUION = AQUI / "guion_hallazgos.md"
GUION.write_text(f"""# Los cuatro hallazgos, para contarlos en voz alta

> Generado por `docs/presentacion/generar_ppt.py` junto con la lámina de la
> auditoría: las cifras salen de los mismos artefactos que la presentación,
> no se escriben aquí a mano. La versión canónica de cada hallazgo vive en
> `app/streamlit_app.py` (lista `AUDITORIA`) y en `INFORME_AUDITORIA.md`.

Un párrafo por hallazgo. Cada uno se puede contar en unos treinta segundos
y responde a lo mismo: qué se leía mal, cómo se notó, en qué quedó.

## 1. El código de faltante, leído como un sueldo

*Nació en los datos de origen.* El INEI codifica «no sabe» como el valor
999999. Ese número entraba al modelo como si fuera un ingreso real de
999.999 soles al mes, y arrastraba consigo toda la regresión: con esos casos
dentro, vivir en zona urbana aparecía asociado a **once soles** más de
ingreso, una cifra sin ningún sentido económico. Al tratar el centinela como
dato faltante, el R² pasó de **{d(_aut['corrida_sucia']['r2'], 3)}** a
**{d(_aut['corrida_limpia']['r2'], 3)}** y todos los coeficientes
recuperaron su signo y su magnitud esperables. Es el hallazgo más grave de
los cuatro y el único que no nació en una decisión nuestra: le habría pasado
a cualquiera que use estos microdatos sin leer el diccionario de variables.

## 2. La rejilla de hiperparámetros estaba acotada

*Decisión propia.* Al buscar los hiperparámetros del modelo desplegado, los
tres ganadores quedaron en el borde del rango que se había probado — la
señal clásica de que el óptimo estaba fuera de la rejilla y nunca se llegó a
mirar. Se amplió el rango y se volvió a buscar: el error baja de
**S/ {REJILLA['vieja']}** a **S/ {REJILLA['nueva']}**, y los tres
hiperparámetros quedan ya en el interior. La mejora es sistemática: gana en
los cinco pliegues de la validación cruzada. Y aun así **no se promovió**.
El motivo no es que la diferencia sea ruido, porque no lo es: es que
{REJILLA['mejora_pct']} % de mejora no justifica regenerar el artefacto
desplegado, revalidar el factor de smearing y rehacer todo el precómputo de
la interfaz. Separar «es mejor» de «vale la pena cambiarlo» es parte del
método, y el hallazgo queda documentado en vez de barrido.

## 3. Una cifra del INEI con la etiqueta equivocada

*Decisión propia.* Se publicaba que el gradiente por tamaño de empresa
«replica el patrón oficial del INEI, con {EXTERNA['inei_1_10']} % de
informalidad en microempresas». La cifra es real y es del INEI, pero
corresponde a su tramo de **1 a 10 trabajadores**, que no es la categoría
«Hasta 20» que usa este proyecto. Al resumir se perdió el tramo, y el lector
terminaba mapeando el {EXTERNA['inei_1_10']} % a una categoría nuestra cuyo
valor propio es **{d(TASA_HASTA20, 1)} %**. No era un dato inventado: era
una comparación mal etiquetada, que es más difícil de detectar precisamente
porque cada mitad, por separado, es correcta. La app ya no la escribe a
mano: la calcula en el precómputo.

## 4. Tres afirmaciones distintas sobre el mismo R²

*Decisión propia y de documentación.* Sobre cuál es el R² esperable en una
ecuación de ingresos circulaban por el proyecto tres versiones a la vez —
«0,4–0,5», «rara vez supera 0,4» y «ningún R² supera 0,5» — repartidas en
cuatro sitios distintos, cada una escrita en un momento diferente y ninguna
consciente de las otras. Al ir a verificarlas apareció el segundo problema:
las dos fuentes que se citaban para sostenerlas, Lemieux (2006) y Heckman,
Lochner y Todd (2006), **no reportan ningún R²**, así que no se las podía
citar para eso. Hoy la afirmación se define **una sola vez**, se apoya en
los cuadros de Mincer y de Card, y dice explícitamente qué parte es lectura
nuestra. La solución de fondo es esa: una cifra, un lugar donde vive.

---

**La línea de método.** Los problemas de origen se corrigen y se documentan;
los propios se corrigen y se aprende de ellos; los de cita se verifican
yendo al texto completo. Ninguno se borra: un hallazgo corregido en silencio
es un hallazgo desperdiciado. Y los dos primeros tenían la misma raíz —
cifras escritas a mano que nadie vuelve a comprobar — así que el arreglo
estructural fue sacarlas del texto y calcularlas en el precómputo.
""", encoding="utf-8")

# --- Reorden según la pauta del docente -----------------------------------
# Las láminas se construyeron en orden de código (1 carátula · 2 datos ·
# 3 variables · 4 push · 5 precómputo · 6 versiones · 7 formulario ·
# 8 tests · 9 torneo · 10 importancia · 11 local=nube · 12 auditoría ·
# 13 prácticas · 14 cierre) y se entregan en el de la pauta: la estadística
# sube antes del bloque de despliegue. Reordenar la lista de diapositivas
# es mover el elemento XML: re-append de un hijo existente lo desplaza.
ORDEN_FINAL = [1, 2, 3, 9, 10, 4, 5, 6, 7, 8, 11, 12, 13, 14]
assert sorted(ORDEN_FINAL) == list(range(1, len(prs.slides._sldIdLst) + 1)), \
    "ORDEN_FINAL no es una permutación de las láminas construidas"
_lst = prs.slides._sldIdLst
_ids = list(_lst)
for _k in ORDEN_FINAL:
    _lst.append(_ids[_k - 1])
TITULOS = [TITULOS[_k - 1] for _k in ORDEN_FINAL]

# --- Doble salida: ENTREGA (fuentes) y EXPO (guion completo) ---------------
SALIDA = AQUI / "ENAHO_exposicion.pptx"
SALIDA_EXPO = AQUI / "ENAHO_exposicion_EXPO.pptx"
_ARG = sys.argv[1] if len(sys.argv) > 1 else ""


def _volcar_notas(expo: bool):
    for _s, _texto, _fuente in NOTAS_REG:
        _s.notes_slide.notes_text_frame.text = (
            _texto if expo else "Fuente: " + _fuente)


if _ARG != "--solo-expo":
    _volcar_notas(expo=False)
    prs.save(SALIDA)
    print(f"{SALIDA.relative_to(RAIZ)} · {len(prs.slides._sldIdLst)} láminas "
          f"· {SALIDA.stat().st_size / 1e6:.2f} MB · ENTREGA (notas mínimas)")
if _ARG != "--solo-entrega":
    _volcar_notas(expo=True)
    prs.save(SALIDA_EXPO)
    print(f"{SALIDA_EXPO.relative_to(RAIZ)} · mismas láminas · EXPO (guion "
          f"del orador)")
print(f"{GUION.relative_to(RAIZ)} · guion de los cuatro hallazgos\n")
print("Títulos-afirmación del mazo (orden de la pauta):")
for k, t in enumerate(TITULOS, 1):
    print(f"  {k:>2}. {t}")
print("\nVerifícalas con:  .venv/Scripts/python.exe "
      "docs/presentacion/verificar_ppt.py [ruta.pptx]")
